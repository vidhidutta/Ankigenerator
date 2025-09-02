import argparse
import os
import yaml
from dotenv import load_dotenv
from pptx import Presentation
from openai import OpenAI
from fpdf import FPDF
import base64
import requests
import re
import time
from sklearn.metrics.pairwise import cosine_similarity
import json
from typing import List, Dict, Tuple, Optional, Any
from dataclasses import dataclass
from difflib import SequenceMatcher
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
import nltk
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.corpus import stopwords
import logging
import genanki
import urllib.request
import shutil
import csv
import subprocess
import tempfile
from ankigenerator.core.image_occlusion import batch_generate_image_occlusion_flashcards
from ankigenerator.core.image_occlusion import cleanup_files
import settings
from anki_models import IOE_MODEL

# Import audio processing (use simplified version for testing)
try:
    from audio_processor import AudioProcessor, AudioMetadata
except ImportError:
    from audio_processor_simple import SimpleAudioProcessor as AudioProcessor, AudioMetadata

from audio_types import AudioBundle

def summarize_card(card):
    if isinstance(card, dict):
        return f"[Card type: {card.get('type', 'unknown')}]"
    else:
        return f"[Card class: {type(card).__name__}, Question: {getattr(card, 'question', '')[:50]}...]"

# Import semantic processing
from semantic_processor import SemanticProcessor

# =====================
# Cloze Utility Helper
# =====================


def clean_cloze_text(text: str) -> str:
    # 1) Iteratively remove any outer {{cN::…}} wrapper:
    outer = re.compile(r'^{{c\d+::(.+)}}$')
    while True:
        m = outer.match(text)
        if not m:
            break
        text = m.group(1)
    # 2) Convert all remaining {{cN::…}} to [text]
    return re.sub(r'{{c\d+::(.*?)}}', r'[\1]', text)

# =====================
# Configuration Section
# =====================
load_dotenv()
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')

# Validate OpenAI API key
if not OPENAI_API_KEY:
    print("❌ Error: OPENAI_API_KEY not found in environment variables")
    print("Please set your OpenAI API key in the .env file:")
    print("OPENAI_API_KEY=your_api_key_here")
    print("Or set it as an environment variable:")
    print("export OPENAI_API_KEY=your_api_key_here")
    # Don't exit here - let the application handle it gracefully
    OPENAI_API_KEY = None

# Load config from YAML
with open('config.yaml', 'r') as f:
    config = yaml.safe_load(f)

PROMPT_TEMPLATE = config.get('prompt_template', """You are an expert medical educator creating Anki flashcards for medical students.

Create concise, focused flashcards that test specific knowledge points. Follow these guidelines:

**Question Style:**
- Keep questions short and direct (1-2 lines max)
- Focus on one specific concept per card
- Use clear, medical terminology
- Avoid verbose explanations in questions

**Answer Style:**
- Provide direct, concise answers (1-3 lines max)
- Include key medical terms and definitions
- For lists, use bullet points or numbered items
- Avoid lengthy explanations unless absolutely necessary

**Examples of Good Questions:**
- "What condition is known for reduced DLCO due to Kco?"
- "What are two concerns for using DLCO for ILD patients?"
- "What stage of ILD is DLCO most important for?"

**Examples of Good Answers:**
- "Interstitial Lung Disease (ILD)"
- "• Falsely reduced in individuals who fail to inspire to TLC\n• Significant variation"
- "Early stage"

**Content Focus:**
- Test key medical concepts, mechanisms, and clinical applications
- Include important clinical signs, symptoms, and diagnostic criteria
- Cover drug mechanisms, side effects, and contraindications
- Test understanding of physiological processes and pathological changes

Generate {flashcard_type} flashcards from the following content:

{batch_text}

Format each flashcard as:
Question: [concise question]
Answer: [direct answer]

Create concise answers that are {cloze}.""")
MODEL_NAME = config.get('model_name', 'gpt-4o')
MAX_TOKENS = config.get('max_tokens', 2000)
TEMPERATURE = config.get('temperature', 0.3)

# User preferences
CATEGORY = config.get('category', 'Other')
EXAM = config.get('exam', 'Other')
ORGANISATION = config.get('organisation', {})
FEATURES = config.get('features', {})
FLASHCARD_TYPE = config.get('flashcard_type', {})
ANSWER_FORMAT = config.get('answer_format', 'best')
CLOZE = config.get('cloze', 'dont_mind')

# Convert flashcard type dictionary to string
def get_flashcard_type_string(flashcard_type_dict):
    """Convert flashcard type dictionary to a meaningful string"""
    if not flashcard_type_dict:
        return "basic"
    
    enabled_types = []
    if flashcard_type_dict.get('level_1', False):
        enabled_types.append("Level 1")
    if flashcard_type_dict.get('level_2', False):
        enabled_types.append("Level 2")
    if flashcard_type_dict.get('both', False):
        enabled_types.append("both levels")
    
    if not enabled_types:
        return "basic"
    elif len(enabled_types) == 1:
        return enabled_types[0]
    else:
        return " and ".join(enabled_types)

FLASHCARD_TYPE_STRING = get_flashcard_type_string(FLASHCARD_TYPE)

# Semantic processing configuration
SEMANTIC_CONFIG = config.get('semantic_processing', {
    'enabled': True,
    'chunk_size': 500,
    'overlap': 50,
    'similarity_threshold': 0.7,
    'embedding_model': 'tfidf'
})

# Initialize semantic processor
semantic_processor = None
if SEMANTIC_CONFIG.get('enabled', True):
    try:
        semantic_processor = SemanticProcessor(
            model_name=SEMANTIC_CONFIG.get('embedding_model', 'tfidf')
        )
        print("✅ Semantic processing initialized successfully (TF-IDF mode)")
    except Exception as e:
        print(f"⚠️ Warning: Could not initialize semantic processing: {e}")
        print("Falling back to basic processing mode")
        semantic_processor = None

# =====================
# Utility Functions
# =====================

def parse_args():
    parser = argparse.ArgumentParser(description='Convert PowerPoint slides to Anki flashcards and extra materials.')
    parser.add_argument('pptx_path', help='Path to the PowerPoint (.pptx) file')
    parser.add_argument('--output', default='flashcards.csv', help='Output CSV file for Anki import')
    parser.add_argument('--notes', default='lecture_notes.pdf', help='Output PDF file for extra materials')
    parser.add_argument('--use_cloze', action='store_true', help='Use cloze cards in generation')
    return parser.parse_args()


def extract_text_from_pptx(pptx_path):
    """Extract text and speaker notes from all slides in a PowerPoint file."""
    try:
        prs = Presentation(pptx_path)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_content = [f"Slide {i+1}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            slide_text = "\n".join(slide_content) if len(slide_content) > 1 else ""
            # Extract speaker notes if present
            notes_text = ""
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
                if notes_frame and notes_frame.text:
                    notes_text = notes_frame.text.strip()
            slides_data.append({
                "slide_text": slide_text,
                "notes_text": notes_text
            })
        print(f"Extracted text and notes from {len(slides_data)} slides")
        return slides_data
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        return []

def convert_pptx_to_slide_pngs(pptx_path: str, output_dir: str) -> List[List[str]]:
    """
    Convert PowerPoint slides to individual PNG images using LibreOffice headless.
    Returns a list of lists, where each inner list contains PNG paths for that slide.
    """
    # Pre-check for LibreOffice CLI
    if shutil.which("libreoffice") is None:
        print("[WARN] LibreOffice not found – slide conversion may fail")
        return []
    
    try:
        # Create a temporary directory for the conversion
        with tempfile.TemporaryDirectory() as temp_dir:
            # Convert PPTX to PDF using LibreOffice headless
            pdf_path = os.path.join(temp_dir, "slides.pdf")
            cmd = [
                "libreoffice", "--headless", "--convert-to", "pdf", 
                "--outdir", temp_dir, pptx_path
            ]
            
            print(f"[DEBUG] Converting PPTX to PDF: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode != 0:
                print(f"[WARN] LibreOffice conversion failed: {result.stderr}")
                return []
            
            # Convert PDF to individual PNG images using pdf2image
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_path, dpi=150)
                
                slide_images = []
                for i, image in enumerate(images):
                    # Save each slide as PNG
                    slide_filename = f"slide{i+1}_full.png"
                    slide_path = os.path.join(output_dir, slide_filename)
                    image.save(slide_path, "PNG")
                    
                    # Add to slide_images list (one image per slide)
                    slide_images.append([slide_path])
                    print(f"[DEBUG] Converted slide {i+1} to: {slide_path}")
                
                print(f"[DEBUG] Successfully converted {len(slide_images)} slides to PNG")
                return slide_images
                
            except ImportError:
                print("[WARN] pdf2image not available. Install with: pip install pdf2image")
                return []
            except Exception as e:
                print(f"[WARN] PDF to PNG conversion failed: {e}")
                return []
                
    except Exception as e:
        print(f"[WARN] Slide conversion failed: {e}")
        return []

def extract_embedded_images_from_pptx(pptx_path: str, output_dir: str) -> List[List[str]]:
    """
    Extract only embedded images from PowerPoint (original function logic).
    Returns a list of lists, where each inner list contains image paths for that slide.
    """
    prs = Presentation(pptx_path)
    slide_images: List[List[str]] = []

    for i, slide in enumerate(prs.slides):
        images_for_slide: List[str] = []
        print(f"[DEBUG] Processing slide {i+1} with {len(slide.shapes)} shapes")

        for j, shp in enumerate(slide.shapes):
            shape_type = getattr(shp, "shape_type", None)
            has_image_attr = hasattr(shp, "image")
            print(f"[DEBUG] Slide {i+1}, Shape {j+1}: type={shape_type}, has_image={has_image_attr}")
            
            # Only process shapes that *have* an image attribute to avoid
            # AttributeError on non-picture shapes (e.g. GraphicFrame)
            if shape_type == 13 and has_image_attr:
                try:
                    pic = shp.image  # type: ignore[attr-defined]
                    image_bytes = pic.blob
                    ext = pic.ext
                    
                    # Skip unsupported image formats
                    if ext.lower() in ['wmf', 'emf']:
                        print(f"⚠️  Skipping unsupported image format on slide {i+1}: {ext}")
                        continue
                    
                    img_name = f"slide{i+1}_img{j+1}.{ext}"
                    img_path = os.path.join(output_dir, img_name)

                    with open(img_path, "wb") as f:
                        f.write(image_bytes)

                    images_for_slide.append(img_path)
                except Exception as e:
                    # Log and continue on any unexpected problem with this picture
                    print(f"⚠️  Skipped image on slide {i+1}: {e}")

        # Always append – even if the list is empty – so indices stay aligned
        slide_images.append(images_for_slide)

    print(
        f"Extracted embedded images for {len(slide_images)} slides "
        f"(total images: {sum(len(lst) for lst in slide_images)})"
    )

    return slide_images

def extract_images_from_pptx(pptx_path, output_dir="slide_images"):
    """Extract all pictures from a PowerPoint file.

    The function now returns a list that has **exactly one entry per slide** so that
    downstream logic can safely `zip(slide_texts, slide_images)` without running
    out-of-sync.  Each entry is a list (possibly empty) of image paths extracted
    from that slide.
    """
    # Ensure the output directory exists
    os.makedirs(output_dir, exist_ok=True)

    # First, try to convert slides to PNG images (captures all content)
    print("[DEBUG] Converting slides to PNG images...")
    slide_pngs = convert_pptx_to_slide_pngs(pptx_path, output_dir)
    
    # Also extract embedded images as a supplement
    print("[DEBUG] Extracting embedded images...")
    embedded_images = extract_embedded_images_from_pptx(pptx_path, output_dir)
    
    # Combine both approaches: prefer slide PNGs, fall back to embedded images
    combined_images = []
    for i in range(max(len(slide_pngs), len(embedded_images))):
        slide_images = []
        
        # Add slide PNG if available
        if i < len(slide_pngs) and slide_pngs[i]:
            slide_images.extend(slide_pngs[i])
        
        # Add embedded images if available
        if i < len(embedded_images) and embedded_images[i]:
            slide_images.extend(embedded_images[i])
        
        # Filter out any paths that are directories or not files
        valid_images = []
        for img_path in slide_images:
            if os.path.exists(img_path):
                if os.path.isfile(img_path):
                    valid_images.append(img_path)
                else:
                    print(f"[WARN] Skipping invalid image path (not a file): {img_path}")
            else:
                print(f"[WARN] Skipping non-existent image path: {img_path}")
        
        combined_images.append(valid_images)
    
    total_images = sum(len(lst) for lst in combined_images)
    print(f"Extracted images for {len(combined_images)} slides (total images: {total_images})")
    
    return combined_images


def stringify_dict(obj: Any, use_cloze: bool = False) -> str:
    """
    Convert a dict—or any object with attributes—to a readable string.
    If use_cloze is True, wraps each value in {{c1::…}} Anki cloze syntax.
    """
    # If obj isn't dict-like, fall back to its attributes
    if not hasattr(obj, "items"):
        try:
            obj = obj.__dict__
        except AttributeError:
            obj = vars(obj)

    items = []
    for key, value in obj.items():
        text = str(value)
        if use_cloze:
            text = f"{{{{c1::{text}}}}}"
        items.append(f"{key}: {text}")
    return "\n".join(items)


def build_extra_materials_prompt(slide_texts, features):
    sections = []
    if features.get('topic_map'):
        sections.append("a topic map (outline of main lecture themes)")
    if features.get('index'):
        sections.append("an index (list of all flashcards)")
    if features.get('glossary'):
        sections.append("a glossary of key terms (with brief definitions)")
    if features.get('summary_review_sheet'):
        sections.append("a summary review sheet (compressed revision version of the whole lecture)")
    if not sections:
        return None  # No extra materials requested
    prompt = (
        "You are an expert medical educator.\n"
        f"Category: {CATEGORY}\nExam: {EXAM}\n"
        "For the following slides, generate:\n- " + "\n- ".join(sections) +
        "\nFormat each section with a clear heading.\n" + "\n\n".join(slide_texts)
    )
    return prompt


def call_api_for_extra_materials(prompt, api_key, model, max_tokens, temperature):
    """Call OpenAI API for extra materials (glossary, etc.)"""
    api_url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    data = {
        "model": model,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "max_tokens": max_tokens,
        "temperature": temperature
    }
    try:
        response = requests.post(api_url, headers=headers, json=data)
        if response.status_code != 200:
            return f"__API_ERROR__Status {response.status_code}: {response.text}"
        response_json = response.json()
        if "choices" not in response_json or not response_json["choices"]:
            return "__API_ERROR__Invalid response format"
        return response_json["choices"][0]["message"]["content"]
    except Exception as e:
        return f"__API_ERROR__{e}"


def parse_flashcards(ai_response, use_cloze=False, slide_number=0, level=1, quality_controller=None):
    flashcards = []
    if ai_response.startswith("__API_ERROR__"):
        print(f"⚠️ OpenAI API error: {ai_response[len('__API_ERROR__'):].strip()}")
        return [], None
    if quality_controller is None:
        from flashcard_generator import QualityController
        quality_controller = QualityController()
    print(f"DEBUG: Parsing AI response of length {len(ai_response)}")
    print(f"DEBUG: First 200 chars: {ai_response[:200]}")
    # Patterns as before...
    pattern_strict = re.compile(r"Question:\s*(.*?)\s*\|\s*Answer:\s*(.*)")
    pattern_qa = re.compile(r"Question:?\s*(.*?)\s*\nAnswer:?\s*(.*?)(?:\n|$)", re.IGNORECASE)
    pattern_md = re.compile(r"\*\*Question:?\*\*\s*(.*?)\s*\n\*\*Answer:?\*\*\s*(.*?)(?:\n|$)", re.IGNORECASE)
    pattern_numbered = re.compile(r"\d+\.\s*Question:\s*(.*?)\s*\n\s*Answer:\s*(.*?)(?=\n\d+\.|$)", re.DOTALL | re.IGNORECASE)
    pattern_simple_numbered = re.compile(r"\d+\.\s*(.*?)\s*\n\s*(.*?)(?=\n\d+\.|$)", re.DOTALL)
    pattern_oneline = re.compile(r"Question:\s*(.*?)\s*\|\s*Answer:\s*(.*?)(?=\nQuestion:|$)", re.DOTALL)
    matches = list(pattern_strict.finditer(ai_response))
    for match in matches:
        question, answer = match.groups()
        if question.strip() and answer.strip():
            flashcards.append((question.strip(), answer.strip()))
    if not flashcards:
        matches = list(pattern_oneline.finditer(ai_response))
        for match in matches:
            question, answer = match.groups()
            if question.strip() and answer.strip():
                flashcards.append((question.strip(), answer.strip()))
    if not flashcards:
        matches = list(pattern_numbered.finditer(ai_response))
        for match in matches:
            question, answer = match.groups()
            if question.strip() and answer.strip():
                flashcards.append((question.strip(), answer.strip()))
    if not flashcards:
        matches = list(pattern_simple_numbered.finditer(ai_response))
        for match in matches:
            question, answer = match.groups()
            if question.strip() and answer.strip() and is_medical_content(question, answer):
                flashcards.append((question.strip(), answer.strip()))
    if not flashcards:
        matches = list(pattern_md.finditer(ai_response))
        for match in matches:
            question, answer = match.groups()
            if question.strip() and answer.strip() and is_medical_content(question, answer):
                flashcards.append((question.strip(), answer.strip()))
    if not flashcards:
        matches = list(pattern_qa.finditer(ai_response))
        for match in matches:
            question, answer = match.groups()
            if question.strip() and answer.strip() and is_medical_content(question, answer):
                flashcards.append((question.strip(), answer.strip()))
    print(f"DEBUG: Total flashcards parsed: {len(flashcards)}")
    # Convert to Flashcard objects and apply cloze if needed
    flashcard_objs = []
    for q, a in flashcards:
        is_cloze = False
        cloze_text = ""
        if use_cloze:
            print(f"[DEBUG] Checking cloze opportunities for: {q[:50]}...")
            is_cloze, cloze_text = quality_controller.detect_cloze_opportunities(q, a)
            print(f"[DEBUG] Cloze result: {is_cloze}, text: {cloze_text[:100] if cloze_text else 'None'}")
            preview_text = clean_cloze_text(cloze_text)
        else:
            # Use a shorter preview: first sentence or up to 100 chars
            try:
                first_sentence = re.split(r"(?<=[.!?])\s+", q.strip())[0]
            except Exception:
                first_sentence = q.strip()
            preview_text = (first_sentence[:100] + ("…" if len(first_sentence) > 100 else "")) if first_sentence else q
        flashcard_objs.append(Flashcard(
            question=q,
            answer=a,
            level=level,
            slide_number=slide_number,
            card_type="basic",
            is_cloze=is_cloze and use_cloze,
            cloze_text=cloze_text if is_cloze and use_cloze else "",
            preview_text=preview_text,
            raw_cloze_text=cloze_text  # store raw cloze text
        ))
    return flashcard_objs, None


def is_medical_content(question, answer):
    """
    Filter out non-medical content like "setting the scene" slides.
    Returns True if the content appears to be medical/exam-relevant.
    """
    # Convert to lowercase for easier matching
    q_lower = question.lower()
    a_lower = answer.lower()
    
    # Skip obvious non-medical content
    skip_patterns = [
        # Generic lecture framework questions
        r"what are the key questions",
        r"what does the.*question focus on",
        r"how does the.*question contribute",
        r"when should the.*question be applied",
        r"what is the purpose of asking",
        
        # Generic analysis frameworks
        r"what\? how\? why\? when\?",
        r"key questions to consider",
        r"medical analysis",
        r"medical decision-making",
        
        # Lecture navigation content
        r"title slide",
        r"contents slide",
        r"lecture outline",
        r"learning objectives",
        r"introduction",
        
        # Too vague or generic
        r"what is.*important",
        r"why is.*important",
        r"how does.*work",
        r"what happens when",
        
        # Lecturer-specific content
        r"lecturer",
        r"professor",
        r"dr\.",
        r"dr ",
        
        # Context-dependent language (references to slides, diagrams, etc.)
        r"mentioned in.*slide",
        r"described in.*slide",
        r"listed on.*slide",
        r"as shown in",
        r"as seen in",
        r"in the slide",
        r"from the slide",
        r"on the slide",
        r"the slide",
        r"this slide",
        r"that slide",
        r"the diagram",
        r"this diagram",
        r"that diagram",
        r"the chart",
        r"this chart",
        r"that chart",
        r"the graph",
        r"this graph",
        r"that graph",
        r"the image",
        r"this image",
        r"that image",
        r"the figure",
        r"this figure",
        r"that figure",
    ]
    
    for pattern in skip_patterns:
        if re.search(pattern, q_lower) or re.search(pattern, a_lower):
            return False
    
    # Must contain medical/clinical terms to be considered relevant
    medical_terms = [
        r"drug", r"medication", r"therapy", r"treatment", r"diagnosis", r"symptom",
        r"disease", r"condition", r"syndrome", r"mechanism", r"receptor", r"enzyme",
        r"protein", r"cell", r"tissue", r"organ", r"system", r"function", r"action",
        r"effect", r"side effect", r"adverse", r"contraindication", r"indication",
        r"dose", r"dosage", r"administration", r"metabolism", r"excretion",
        r"pharmacokinetics", r"pharmacodynamics", r"interaction", r"toxicity",
        r"efficacy", r"potency", r"selectivity", r"specificity", r"affinity",
        r"agonist", r"antagonist", r"inhibitor", r"activator", r"modulator",
        r"pathway", r"cascade", r"signal", r"transduction", r"regulation",
        r"homeostasis", r"balance", r"equilibrium", r"threshold", r"baseline",
        r"normal", r"abnormal", r"pathological", r"physiological", r"clinical",
        r"patient", r"case", r"presentation", r"history", r"examination",
        r"investigation", r"test", r"result", r"finding", r"observation",
        r"assessment", r"evaluation", r"management", r"care", r"monitoring",
        r"follow-up", r"outcome", r"prognosis", r"complication", r"risk",
        r"factor", r"etiology", r"pathogenesis", r"pathophysiology", r"anatomy",
        r"physiology", r"biochemistry", r"molecular", r"genetic", r"immunology",
        r"microbiology", r"infection", r"bacteria", r"virus", r"fungus",
        r"parasite", r"antibiotic", r"antiviral", r"antifungal", r"vaccine",
        r"immunization", r"allergy", r"hypersensitivity", r"autoimmune",
        r"inflammation", r"injury", r"trauma", r"surgery", r"procedure",
        r"technique", r"method", r"approach", r"strategy", r"protocol",
        r"guideline", r"recommendation", r"standard", r"practice", r"policy"
    ]
    
    # Check if question or answer contains medical terms
    for term in medical_terms:
        if re.search(term, q_lower) or re.search(term, a_lower):
            return True
    
    # If no medical terms found, it's likely not medical content
    return False


def export_flashcards_to_apkg(flashcards, output_path='flashcards.apkg', pptx_filename=None, config_topic=None):
    import base64
    import os

    # Detect topic from filename or config
    if config_topic:
        topic = config_topic
    elif pptx_filename:
        topic = os.path.splitext(os.path.basename(pptx_filename))[0].replace('_', ' ').replace('-', ' ').title()
    else:
        topic = "Lecture"

    # Define Anki note models
    BASIC_MODEL = genanki.Model(
        1607392319,
        'Basic (My Generated Deck)',
        fields=[
            {'name': 'Question'},
            {'name': 'Answer'},
            {'name': 'Audio'},
        ],
        templates=[
            {
                'name': 'Card 1',
                'qfmt': '{{Question}}<br>{{Audio}}',
                'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}<br>{{Audio}}',
            },
        ],
    )

    CLOZE_MODEL = genanki.Model(
        998877661,
        'Cloze (My Generated Deck)',
        fields=[
            {'name': 'Text'},
            {'name': 'Audio'},
        ],
        templates=[
            {
                'name': 'Cloze Card',
                'qfmt': '{{cloze:Text}}<br>{{Audio}}',
                'afmt': '{{cloze:Text}}<br>{{Audio}}',
            },
        ],
        model_type=genanki.Model.CLOZE,
    )

    # Use a stable but fairly unique deck ID – hash the topic string
    deck_id = abs(hash(topic)) % 2147483647  # Anki requires 32-bit signed int
    deck = genanki.Deck(deck_id, topic)

    print("[DEBUG] export_flashcards_to_apkg() called")
    print(f"[DEBUG] Number of flashcards to export: {len(flashcards)}")
    
    # Better debug information for flashcard types
    type_counts = {}
    for card in flashcards:
        if isinstance(card, dict):
            card_type = card.get('type', 'unknown')
        elif hasattr(card, 'card_type'):
            card_type = card.card_type or 'text'
        elif hasattr(card, 'question'):  # Flashcard object
            card_type = 'text'
        else:
            card_type = 'unknown'
        
        type_counts[card_type] = type_counts.get(card_type, 0) + 1
    
    print(f"[DEBUG] Flashcard type counts: {type_counts}")

    print("\n===== Exporting Flashcards =====")
    media = []  # Collect image paths for the package

    # Lazy import of Flashcard class to avoid cyclic issues
    try:
        from flashcard_generator import Flashcard  # type: ignore
    except Exception:
        Flashcard = None  # Fallback – isinstance checks will fail gracefully

    for entry in flashcards:
        # -----------------------------
        # IMAGE-OCCLUSION DICT HANDLING
        # -----------------------------
        if isinstance(entry, dict):
            if 'question_image_path' in entry and 'answer_image_path' in entry:
                # Validate image paths are not directories
                if os.path.isdir(entry.get('question_image_path', '')):
                    print(f"[WARN] Skipping invalid flashcard: question image path is directory - {entry}")
                    continue
                if os.path.isdir(entry.get('answer_image_path', '')):
                    print(f"[WARN] Skipping invalid flashcard: answer image path is directory - {entry}")
                    continue
                
                # Handle both file paths and Base64 encoded images
                if 'question_image_base64' in entry and 'answer_image_base64' in entry:
                    # Use Base64 encoded images
                    qfile = f"data:image/png;base64,{entry['question_image_base64']}"
                    afile = f"data:image/png;base64,{entry['answer_image_base64']}"
                    print(f"[DEBUG] Using Base64 encoded images for image occlusion card")
                else:
                    # Use file paths (legacy support)
                    qfile = os.path.basename(entry['question_image_path'])
                    afile = os.path.basename(entry['answer_image_path'])
                    # Track media files
                    media.extend([entry['question_image_path'], entry['answer_image_path']])
                    print(f"[DEBUG] Using file paths for image occlusion card")

                # Audio embedding for dict entries if provided
                audio_field = ""
                dict_audio_files = entry.get('audio_files') or []
                for apath in dict_audio_files[:2]:
                    if apath and os.path.isfile(apath):
                        media.append(apath)
                        audio_field += f"[sound:{os.path.basename(apath)}] "

                # Alt text
                safe_alt = (entry.get('answer_text') or '').replace('"', '&quot;')
                front_html = f"<img src='{qfile}' alt='Covered region'>"
                back_html = f"<img src='{afile}' alt='{safe_alt}'>"

                note = genanki.Note(
                    model=IOE_MODEL,
                    fields=[front_html, back_html],
                )
                # If optional answer_text/rationale present, append below back image
                ans_txt = (entry.get('answer_text') or '').strip()
                rationale = (entry.get('rationale') or '').strip()
                # Toggle from config
                show_rat = bool(config.get('image_understanding', {}).get('show_rationale_on_back', True))
                if ans_txt or (rationale and show_rat):
                    back_extra = "".join([
                        f"<div style='margin-top:8px;font-size:12px;color:#333;'><b>Answer:</b> {ans_txt}</div>" if ans_txt else "",
                        (f"<div style='margin-top:4px;font-size:11px;color:#666;'>{rationale}</div>" if rationale else "") if show_rat else "",
                    ])
                    note.fields[1] = note.fields[1] + back_extra
                deck.add_note(note)
            else:
                print(f"[WARN] Skipping dict entry missing expected keys: {entry}")
            continue

        # -----------------------------
        # TEXT FLASHCARD OBJECTS
        # -----------------------------
        if Flashcard is not None and isinstance(entry, Flashcard):
            # Determine note model based on cloze status
            if getattr(entry, 'is_cloze', False):
                text_field = entry.cloze_text or f"{entry.question} – {entry.answer}"
                # Prepare audio field
                audio_field = ""
                if getattr(entry, 'audio_metadata', None):
                    audio_files = getattr(entry.audio_metadata, 'audio_files', []) or []
                    for apath in audio_files[:2]:
                        if apath and os.path.isfile(apath):
                            media.append(apath)
                            audio_field += f"[sound:{os.path.basename(apath)}] "
                note = genanki.Note(model=CLOZE_MODEL, fields=[text_field, audio_field.strip()])
            else:
                # Prepare audio field
                audio_field = ""
                if getattr(entry, 'audio_metadata', None):
                    audio_files = getattr(entry.audio_metadata, 'audio_files', []) or []
                    for apath in audio_files[:2]:
                        if apath and os.path.isfile(apath):
                            media.append(apath)
                            audio_field += f"[sound:{os.path.basename(apath)}] "
                note = genanki.Note(model=BASIC_MODEL, fields=[entry.question, entry.answer, audio_field.strip()])

            # Comprehensive path validation for all potential path attributes
            path_attributes = ['image_path', 'question_image_path', 'answer_image_path']
            for attr in path_attributes:
                path_value = getattr(entry, attr, None)
                if path_value and os.path.exists(path_value):
                    if os.path.isdir(path_value):
                        print(f"[WARN] Skipping invalid flashcard: {attr} is directory - {entry}")
                        continue

            # If the Flashcard references an image_path, embed it for context
            img_path = getattr(entry, 'image_path', None)
            if img_path and os.path.exists(img_path):
                # Validate that image_path is actually a file, not a directory
                if os.path.isfile(img_path):
                    fname = os.path.basename(img_path)
                    media.append(img_path)

                    # Append image below the answer/cloze content.
                    # For Basic model (2+ fields) -> use Answer field (index 1)
                    if len(note.fields) >= 2:
                        note.fields[1] += f"<br><img src='{fname}'>"
                    else:
                        # Cloze card – single field; append image to same field
                        note.fields[0] += f"<br><img src='{fname}'>"
                else:
                    print(f"[WARN] Skipping image_path - is not a file: {img_path}")

            deck.add_note(note)
            continue

        # -----------------------------
        # UNKNOWN ENTRY TYPE
        # -----------------------------
        print(f"[WARN] Unrecognized flashcard entry – skipped: {entry}")

    files_to_cleanup: list[str] = []
    try:
        print(f"[DEBUG] Media files to include: {media}")
        print(f"[DEBUG] Number of media files: {len(media)}")
        
        # Filter out invalid media files (directories, non-existent files)
        valid_media = []
        for i, media_file in enumerate(media):
            if os.path.isdir(media_file):
                print(f"[WARN] Skipping media file {i} - is a directory: {media_file}")
                continue
            elif not os.path.isfile(media_file):
                print(f"[WARN] Skipping media file {i} - not a file: {media_file}")
                continue
            else:
                valid_media.append(media_file)
        
        print(f"[DEBUG] Valid media files: {len(valid_media)} out of {len(media)}")
        
        # Validate that all [sound:...] references exist in valid_media
        referenced = set()
        try:
            for n in deck.notes:
                for field in n.fields:
                    if isinstance(field, str) and '[sound:' in field:
                        import re
                        for m in re.finditer(r"\[sound:([^\]]+)\]", field):
                            referenced.add(m.group(1))
        except Exception:
            pass
        missing = []
        vm_basenames = {os.path.basename(p) for p in valid_media}
        for r in referenced:
            if r not in vm_basenames:
                missing.append(r)
        if missing:
            raise RuntimeError(f"Missing audio media files in package: {missing}")
        
        genanki.Package(deck, media_files=valid_media).write_to_file(output_path)
    except Exception as e:
        print(f"[ERROR] Full error: {e}")
        cleanup_files(files_to_cleanup)
        raise
    print("[DEBUG] Finished adding flashcards, writing to file...")
    print(f"[DEBUG] Successfully wrote APKG file to: {output_path}")
    print(f"[SUCCESS] APKG exported to {output_path}")
    
    return output_path



def save_text_to_pdf(text, filename):
    # Use a Unicode font for full character support
    font_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")
    font_path = os.path.join(font_dir, "DejaVuSans.ttf")
    if not os.path.exists(font_path):
        raise FileNotFoundError(
            f"DejaVuSans.ttf not found in {font_dir}. Please download it and place it in the fonts/ directory."
        )
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_font("DejaVu", "", font_path, uni=True)
    pdf.set_font("DejaVu", size=12)
    for line in text.split('\n'):
        pdf.cell(0, 10, line, ln=True)
    pdf.output(filename)


def generate_multimodal_flashcards_http(slide_texts, slide_images, api_key, model, max_tokens, temperature, use_cloze=False):
    import os, base64
    import requests
    import time
    
    # Validate API key
    if not api_key:
        print("❌ Error: OpenAI API key is required for flashcard generation")
        print("Please set your OpenAI API key in the .env file or environment variables")
        return []
    
    api_url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    all_flashcards = []

    for slide_index in range(len(slide_texts)):
        # Per-slide initializations
        image_paths = slide_images[slide_index] if slide_index < len(slide_images) else []
        last_img_path = None
        content_blocks = [{"type": "text", "text": slide_texts[slide_index]}]

        # Embed each image and remember only the last path
        for img_path in image_paths:
            last_img_path = img_path
            try:
                if os.path.isfile(img_path):
                    with open(img_path, "rb") as f:
                        ext = os.path.splitext(img_path)[1][1:]
                        b64 = base64.b64encode(f.read()).decode("utf-8")
                    content_blocks.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/{ext};base64,{b64}"}
                    })
                else:
                    print(f"Warning: {img_path} is not a file (might be a directory)")
            except Exception as e:
                print(f"Warning: Could not process image {img_path}: {e}")

        # Create prompt for this slide
        prompt = PROMPT_TEMPLATE.format(
            flashcard_type=FLASHCARD_TYPE_STRING,
            cloze=CLOZE,
            batch_text=slide_texts[slide_index]
        )

        # Prepare API request
        data = {
            "model": model,
            "messages": [
                {"role": "user", "content": [{"type": "text", "text": prompt}] + content_blocks[1:]}
            ],
            "max_tokens": max_tokens,
            "temperature": temperature
        }

        try:
            response = requests.post(api_url, headers=headers, json=data, timeout=30)
            response.raise_for_status()
            response_json = response.json()
            
            if "choices" not in response_json or not response_json["choices"]:
                print(f"Error: No choices in response for slide {slide_index + 1}")
                continue
                
            ai_content = response_json["choices"][0]["message"]["content"]
            
            # Parse flashcards from the response
            flashcard_objs, _ = parse_flashcards(
                ai_content,
                use_cloze=use_cloze,
                slide_number=slide_index,
                level=1
            )
            
            # Attach image_path information to each generated Flashcard
            for fc in flashcard_objs:
                if last_img_path and os.path.isfile(last_img_path):
                    fc.image_path = last_img_path
                elif last_img_path:
                    print(f"[WARN] Skipping invalid image path (not a file): {last_img_path}")
            
            all_flashcards.extend(flashcard_objs)
            print(f"Generated {len(flashcard_objs)} flashcards from slide {slide_index + 1}")
            
            time.sleep(0.1)  # Rate limiting
            
        except Exception as e:
            print(f"Error generating flashcards for slide {slide_index + 1}: {e}")
            continue

    return all_flashcards




def filter_slides(slide_texts, slide_images):
    """
    Filter out slides that are likely to be empty or contain only navigation content.
    Returns a list of slide texts that should be processed and their indices.
    """
    filtered_slides = []
    kept_indices = []
    for i, (slide_text, images) in enumerate(zip(slide_texts, slide_images)):
        # Skip slides that are likely to be empty or navigation
        if should_skip_slide(slide_text, images):
            print(f"[DEBUG] Skipping slide {i+1}: no text, no images")
            continue
        print(f"[DEBUG] Keeping slide {i+1}: has text or images")
        filtered_slides.append(slide_text)
        kept_indices.append(i)
    return filtered_slides, kept_indices


def should_skip_slide(slide_text, images):
    """
    Determine if a slide should be skipped based on its content.
    Returns True if the slide should be skipped.
    """
    # Handle both string and dictionary slide_text
    if isinstance(slide_text, dict):
        slide_text_content = slide_text.get('content', '')
    else:
        slide_text_content = str(slide_text)
    
    text_lower = slide_text_content.lower()
    
    # Skip patterns for slides that should be ignored
    skip_patterns = [
        # Empty or nearly empty slides
        r"^slide \d+:\s*$",
        r"^slide \d+:\s*\n\s*$",
        
        # Title slides
        r"^slide \d+:\s*[a-z\s]+\s*$",  # Just a title
        
        # Contents/navigation slides
        r"contents",
        r"outline",
        r"agenda",
        r"learning objectives",
        r"objectives",
        r"introduction",
        r"overview",
        r"summary",
        r"conclusion",
        r"references",
        r"bibliography",
        r"further reading",
        r"questions",
        r"discussion",
        
        # Generic framework slides
        r"what\? how\? why\? when\?",
        r"key questions",
        r"framework",
        r"approach",
        r"methodology",
        
        # Lecturer information
        r"lecturer",
        r"professor",
        r"dr\.",
        r"dr ",
        r"presented by",
        r"by ",
        
        # Too short to be meaningful
        r"^slide \d+:\s*\w+\s*$",  # Just one word
    ]
    
    for pattern in skip_patterns:
        if re.search(pattern, text_lower):
            return True
    
    # If slide has very little content (less than 50 characters excluding "Slide X:") and no images
    slide_match = re.match(r"^slide \d+:", slide_text_content)
    slide_header_length = len(slide_match.group(0)) if slide_match else 0
    content_length = len(slide_text_content) - slide_header_length
    if content_length < 50 and not images:
        return True
    
    return False

# =====================
# Enhanced Flashcard Generation with Semantic Processing
# =====================

def generate_flashcards_from_semantic_chunks(semantic_chunks, slide_images, api_key, model, max_tokens, temperature, progress=None, use_cloze=False, question_style="Word for word"):
    # Validate API key
    if not api_key:
        print("❌ Error: OpenAI API key is required for flashcard generation")
        print("Please set your OpenAI API key in the .env file or environment variables")
        return []
    
    api_url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    all_flashcards = []
    total_chunks = len(semantic_chunks)
    quality_controller = QualityController()
    for i, chunk_data in enumerate(semantic_chunks):
        image_paths = []
        last_img_path = None  # Track last image path for this chunk
        if progress:
            progress_percent = 0.4 + (0.45 * (i / total_chunks))
            progress(progress_percent, desc=f"Processing semantic chunk {i+1}/{total_chunks}...")
        if semantic_processor is not None:
            # Add logging to show what content is causing type mismatch
            print(f"[DEBUG] Processing semantic chunk {i+1}: type={type(chunk_data)}, content={str(chunk_data)[:100]}...")
            
            # Validate input types before calling semantic processor
            if not isinstance(chunk_data, dict):
                print(f"[WARN] chunk_data is not a dict: {type(chunk_data)} - converting to string")
                chunk_data = {'text': str(chunk_data), 'slide_index': 0}
            elif 'text' not in chunk_data:
                print(f"[WARN] chunk_data missing 'text' key: {list(chunk_data.keys())} - using first value")
                # Use the first available value as text
                first_value = next(iter(chunk_data.values()), '')
                chunk_data = {'text': str(first_value), 'slide_index': chunk_data.get('slide_index', 0)}
            
            # Ensure only segment["text"] (or equivalent string field) is passed
            if not isinstance(chunk_data.get('text', ''), str):
                print(f"[WARN] chunk_data['text'] is not a string: {type(chunk_data.get('text'))} - converting")
                chunk_data['text'] = str(chunk_data.get('text', ''))
            
            # Ensure slide_index exists
            if 'slide_index' not in chunk_data:
                print(f"[WARN] chunk_data missing 'slide_index' - adding default value")
                chunk_data['slide_index'] = 0
            
            # Fetch audio context per slide in this chunk if available via global last bundle kept in semantic_processor (none), or pass via chunk_data if present
            audio_context_snippets = []
            try:
                bundle: Optional[AudioBundle] = chunk_data.get('audio_bundle')  # may not exist
                if bundle:
                    slide_numbers = chunk_data.get('slide_numbers') or [chunk_data.get('slide_index', 0)]
                    for slide_num in slide_numbers:
                        wins = [w for w in getattr(bundle, 'slide_windows', []) if w.slide_id == slide_num]
                        if not wins:
                            continue
                        win = wins[0]
                        # Rank segments inside window by emphasis * length * novelty
                        seen_terms = set()
                        scored = []
                        for s in win.segments:
                            length = max(0.0, s.end - s.start)
                            # novelty by new keywords in this segment
                            kws = set(re.findall(r"[A-Za-z0-9_]+", s.text.lower()))
                            new_terms = len([t for t in kws if t not in seen_terms])
                            for t in kws:
                                seen_terms.add(t)
                            score = float(getattr(s, 'emphasis', 0.0)) * length * max(1.0, float(new_terms))
                            scored.append((score, s))
                        scored.sort(key=lambda x: x[0], reverse=True)
                        top = [ss for _, ss in scored[:2]]
                        for ss in top:
                            audio_context_snippets.append(f"[{ss.start:.1f}-{ss.end:.1f}s] {ss.text}")
            except Exception:
                pass
            base_prompt = semantic_processor.build_enhanced_prompt(chunk_data, PROMPT_TEMPLATE)
            if audio_context_snippets:
                enhanced_prompt = base_prompt + "\n\nAudio emphasis snippets (use these to focus and prioritize facts):\n" + "\n".join(f"- {t}" for t in audio_context_snippets)
            else:
                enhanced_prompt = base_prompt
        else:
            # Add question style instructions to the prompt
            style_instructions = ""
            if question_style == "Word for word":
                style_instructions = """
 **Question Style: Word for word**
 - Keep answers in the exact same format as presented in the slide
 - Use the precise terminology and phrasing from the original content
 - Maintain the original structure and order of information
 - For Level 2 questions: Use reasoning questions but keep answers in exact slide format
 """
            elif question_style == "Elaborated":
                style_instructions = """
 **Question Style: Elaborated**
 - Keep answers in the same format as the slide content
 - Add explanations below the original answer to help clarify concepts
 - Provide additional context and clinical significance
 - Format: [Original answer] + [Explanation/context]
 - For Level 1: Keep explanations brief and focused on basic understanding
 - For Level 2: Add deeper clinical reasoning and implications
 """
            elif question_style == "Simplified":
                style_instructions = """
 **Question Style: Simplified**
 - Take complex explanations and convert them to concise, simple phrases
 - Use clear, straightforward language that's easy to understand
 - Break down complex concepts into digestible parts
 - Maintain accuracy while improving clarity
 - For Level 1: Focus on basic definitions and facts
 - For Level 2: Simplify complex reasoning and patterns
 """

            # Ensure chunk_data['text'] is a string with validation
            if not isinstance(chunk_data, dict):
                print(f"[WARN] chunk_data is not a dict in non-semantic mode: {type(chunk_data)}")
                chunk_text = str(chunk_data)
            else:
                chunk_text = chunk_data.get('text', '')
                if isinstance(chunk_text, dict):
                    print(f"[WARN] chunk_data['text'] is a dict: {type(chunk_text)} - converting to string")
                    chunk_text = str(chunk_text)
                elif not isinstance(chunk_text, str):
                    print(f"[WARN] chunk_data['text'] is not a string: {type(chunk_text)} - converting")
                    chunk_text = str(chunk_text)
            
            # Validate input types before calling PROMPT_TEMPLATE.format
            if not isinstance(chunk_text, str):
                print(f"[ERROR] chunk_text is still not a string after conversion: {type(chunk_text)}")
                chunk_text = str(chunk_text)
            
            enhanced_prompt = PROMPT_TEMPLATE.format(
                flashcard_type=FLASHCARD_TYPE_STRING,
                cloze=CLOZE,
                batch_text=chunk_text
            ) + style_instructions
        content_blocks = [
            {"type": "text", "text": enhanced_prompt}
        ]
        slide_index = chunk_data['slide_index']
        if 0 <= slide_index < len(slide_images):
            image_paths = slide_images[slide_index]
            for img_path in image_paths:
                last_img_path = img_path
                try:
                    # Check if the path is actually a file, not a directory
                    if os.path.isfile(img_path):
                        with open(img_path, "rb") as image_file:
                            ext = os.path.splitext(img_path)[1][1:]
                            base64_image = base64.b64encode(image_file.read()).decode("utf-8")
                            content_blocks.append({
                                "type": "image_url",
                                "image_url": {"url": f"data:image/{ext};base64,{base64_image}"}
                            })
                    else:
                        print(f"Warning: {img_path} is not a file (might be a directory)")
                except Exception as e:
                    print(f"Warning: could not read {img_path}: {e}")
        data = {
            "model": model,
            "messages": [
                {"role": "user", "content": content_blocks}
            ],
            "max_tokens": max_tokens,
            "temperature": temperature
        }
        try:
            response = requests.post(api_url, headers=headers, json=data)
            if response.status_code != 200:
                print(f"Error: Received status code {response.status_code} for chunk {i+1}")
                print("Response text:", response.text)
                continue
            try:
                response_json = response.json()
            except Exception as e:
                print(f"Error decoding JSON for chunk {i+1}: {e}")
                print("Raw response:", response.text)
                continue
            if "choices" not in response_json or not response_json["choices"]:
                print(f"Error: 'choices' key missing or empty in API response for chunk {i+1}")
                print("Full response:", response_json)
                continue
            ai_content = response_json["choices"][0]["message"]["content"]
            print(f"\n--- AI RAW RESPONSE FOR SEMANTIC CHUNK {i+1} ---\n", ai_content, "\n----------------------\n")
            # Parse flashcards from the response as Flashcard objects
            flashcard_objs, _ = parse_flashcards(
                ai_content,
                use_cloze=use_cloze,
                slide_number=slide_index,
                level=1,  # Default to level 1; can be improved if level info is available
                quality_controller=quality_controller
            )
            # Attach image_path information to each generated Flashcard
            for fc in flashcard_objs:
                # Only set image_path if it's a verified file
                if last_img_path and os.path.isfile(last_img_path):
                    fc.image_path = last_img_path
                elif last_img_path:
                    print(f"[WARN] Skipping invalid image path (not a file): {last_img_path}")
            print(f"Generated {len(flashcard_objs)} flashcards from semantic chunk {i+1}")
            all_flashcards.extend(flashcard_objs)
            if progress:
                total_generated = len(all_flashcards)
                progress(progress_percent, desc=f"Chunk {i+1}/{total_chunks}: Generated {len(flashcard_objs)} cards (Total: {total_generated})")
            time.sleep(0.1)
        except Exception as e:
            print(f"Error generating flashcards for semantic chunk {i+1}: {e}")
    # After building all_flashcards, try to sort by audio emphasis provenance if available
    try:
        def card_emphasis(card):
            md = getattr(card, 'audio_metadata', None)
            if md and getattr(md, 'audio_files', None):
                # approximate: if we had provenance scores, we'd use them; fall back to emphasis_weight
                return getattr(card, 'emphasis_weight', 0.0)
            return getattr(card, 'emphasis_weight', 0.0)
        all_flashcards.sort(key=card_emphasis, reverse=True)
    except Exception:
        pass
    return all_flashcards

def generate_enhanced_flashcards_with_progress(slide_texts, slide_images, api_key, model, max_tokens, temperature, progress=None, use_cloze=False, question_style="Word for word", audio_bundle=None):
    """
    Generate flashcards using semantic processing with progress tracking
    
    Args:
        slide_texts: List of slide texts
        slide_images: List of image paths for each slide
        api_key: OpenAI API key
        model: Model name
        max_tokens: Max tokens for generation
        temperature: Temperature for generation
        progress: Progress callback function
        use_cloze: Boolean indicating whether to use cloze cards
        question_style: Style of questions ("Word for word", "Elaborated", "Simplified")
        audio_bundle: Optional AudioBundle providing audio segments and slide alignment
        
    Returns:
        Tuple of (flashcards, analysis_data)
    """
    # Initialize audio context from bundle if provided
    slide_audio_map = {}
    content_weights = {}
    
    if audio_bundle:
        try:
            if progress:
                progress(0.05, desc="Processing audio bundle...")
            # Build slide_audio_map from bundle
            for win in getattr(audio_bundle, 'slide_windows', []) or []:
                segs = getattr(win, 'segments', []) or []
                slide_audio_map[win.slide_id] = segs
            # Compute content weights using a lightweight heuristic if available
            try:
                from audio_processor import AudioProcessor
                ap = AudioProcessor(model_name="base")
                content_weights = ap.calculate_content_weights_from_bundle(audio_bundle)
            except Exception:
                # Fallback: weight by avg emphasis
                content_weights = {}
                for win in getattr(audio_bundle, 'slide_windows', []) or []:
                    if win.segments:
                        avg_emphasis = float(np.mean([getattr(s, 'emphasis', 0.5) for s in win.segments]))
                        content_weights[win.slide_id] = avg_emphasis
            # Logging
            all_segments = getattr(audio_bundle, 'segments', []) or []
            print(f"🎵 Audio Analysis:")
            print(f"   • Total audio segments: {len(all_segments)}")
            print(f"   • Slides with audio: {len(slide_audio_map)}")
            if all_segments:
                print(f"   • Average emphasis score: {np.mean([getattr(seg, 'emphasis', 0.5) for seg in all_segments]):.2f}")
        except Exception as e:
            print(f"⚠️ Audio bundle processing failed: {e}")
            print("Continuing without audio context...")
    
    if semantic_processor is None:
        print("⚠️ Semantic processing not available, falling back to basic processing")
        return generate_multimodal_flashcards_http(slide_texts, slide_images, api_key, model, max_tokens, temperature, use_cloze=use_cloze), None
    
    try:
        # Step 1: Create semantic chunks
        if progress:
            progress(0.4, desc="Creating semantic chunks and analyzing content...")
        # Normalize slide_texts to strings (OCR augmentation may store dicts)
        try:
            normalized_slide_texts = []
            for st in slide_texts:
                if isinstance(st, dict):
                    normalized_slide_texts.append(str(st.get('text', '')))
                else:
                    normalized_slide_texts.append(str(st))
        except Exception:
            normalized_slide_texts = [str(st) for st in slide_texts]
        semantic_chunks = semantic_processor.create_semantic_chunks(normalized_slide_texts)
        
        # Attach audio bundle reference to chunks so downstream can use it
        try:
            if audio_bundle:
                for ch in semantic_chunks:
                    if isinstance(ch, dict):
                        ch['audio_bundle'] = audio_bundle
        except Exception:
            pass
        
        # Step 2: Analyze content quality
        analysis_data = semantic_processor.analyze_content_quality(semantic_chunks)
        print(f"📊 Content Analysis:")
        print(f"   • Total chunks: {analysis_data['total_chunks']}")
        print(f"   • Total slides: {analysis_data['total_slides']}")
        print(f"   • Average chunk size: {analysis_data['avg_chunk_size']:.1f} characters")
        print(f"   • Unique key phrases: {analysis_data['unique_key_phrases']}")
        print(f"   • Average group size: {analysis_data['avg_group_size']:.1f}")
        print(f"   • Key phrases: {', '.join(analysis_data['key_phrases'][:5])}")
        # Step 3: Generate flashcards from semantic chunks with audio context
        all_flashcards = generate_flashcards_from_semantic_chunks(
            semantic_chunks, slide_images, api_key, model, max_tokens, temperature, progress, use_cloze, question_style
        )
        # Apply audio weighting if available
        if slide_audio_map:
            all_flashcards = apply_audio_weighting(all_flashcards, slide_audio_map, content_weights)
        print("\nSample of generated flashcards (before export):")
        for card in all_flashcards[:10]:
            print(summarize_card(card))
        return all_flashcards, analysis_data
    except Exception as e:
        error_msg = str(e)
        print(f"Error in enhanced flashcard generation: {error_msg}")
        print(f"Error type: {type(e)}")
        print(f"Error details: {e}")
        import traceback
        print(f"Full traceback: {traceback.format_exc()}")
        print("Falling back to basic processing")
        # Return the error message so the UI can handle it appropriately
        return generate_multimodal_flashcards_http(slide_texts, slide_images, api_key, model, max_tokens, temperature, use_cloze=use_cloze), None

def enhance_chunk_with_audio_context(chunk: Dict[str, Any], slide_audio_map: Dict[int, List], 
                                   content_weights: Dict[int, float], audio_processor: AudioProcessor) -> Dict[str, Any]:
    """
    Enhance semantic chunk with audio context information
    
    Args:
        chunk: Semantic chunk data
        slide_audio_map: Mapping of slides to audio segments
        content_weights: Importance weights for each slide
        audio_processor: AudioProcessor instance
        
    Returns:
        Enhanced chunk with audio context
    """
    enhanced_chunk = chunk.copy()
    
    # Add audio context to the prompt
    audio_context = []
    for slide_num in chunk.get('slide_numbers', []):
        if slide_num in slide_audio_map and slide_audio_map[slide_num]:
            segments = slide_audio_map[slide_num]
            weight = content_weights.get(slide_num, 0.5)
            
            # Extract audio metadata
            audio_metadata = audio_processor.extract_audio_metadata(slide_num, segments)
            
            # Add emphasis information
            if audio_metadata.emphasis_score > 0.7:
                audio_context.append(f"Slide {slide_num}: Lecturer emphasized this content (emphasis score: {audio_metadata.emphasis_score:.2f})")
            elif audio_metadata.emphasis_score < 0.3:
                audio_context.append(f"Slide {slide_num}: Lecturer covered this content briefly (emphasis score: {audio_metadata.emphasis_score:.2f})")
            
            # Add time allocation information
            if audio_metadata.time_allocation > 30:
                audio_context.append(f"Slide {slide_num}: Lecturer spent {audio_metadata.time_allocation:.1f} seconds on this content")
            elif audio_metadata.time_allocation < 10:
                audio_context.append(f"Slide {slide_num}: Lecturer spent only {audio_metadata.time_allocation:.1f} seconds on this content")
    
    if audio_context:
        enhanced_chunk['audio_context'] = "\n".join(audio_context)
        enhanced_chunk['content'] += f"\n\nAudio Context:\n" + "\n".join(audio_context)
    
    return enhanced_chunk

def apply_audio_weighting(flashcards: List['Flashcard'], slide_audio_map: Dict[int, List], 
                         content_weights: Dict[int, float]) -> List['Flashcard']:
    """
    Apply audio-based weighting to flashcards
    
    Args:
        flashcards: List of flashcards
        slide_audio_map: Mapping of slides to audio segments
        content_weights: Importance weights for each slide
        
    Returns:
        List of flashcards with audio weighting applied
    """
    for flashcard in flashcards:
        slide_num = flashcard.slide_number
        if slide_num in content_weights:
            weight = content_weights[slide_num]
            flashcard.emphasis_weight = min(max(weight * 2, 0.5), 2.0)
            if slide_num in slide_audio_map and slide_audio_map[slide_num]:
                segments = slide_audio_map[slide_num]
            total_time = 0.0
            for seg in segments:
                start = getattr(seg, 'start_time', getattr(seg, 'start', None))
                end = getattr(seg, 'end_time', getattr(seg, 'end', None))
                if start is not None and end is not None:
                    total_time += max(end - start, 0.0)
                flashcard.time_allocation = total_time
    return flashcards

def remove_duplicate_flashcards(flashcards, similarity_threshold=0.8):
    """
    Remove duplicate or very similar flashcards using semantic similarity
    
    Args:
        flashcards: List of Flashcard objects or (question, answer) tuples
        similarity_threshold: Threshold for considering flashcards similar
        
    Returns:
        List of deduplicated flashcards
    """
    if not flashcards or semantic_processor is None:
        return flashcards
    
    try:
        # Handle both Flashcard objects and tuples
        if isinstance(flashcards[0], Flashcard):
            # Extract questions and answers from Flashcard objects
            questions = [fc.question for fc in flashcards]
            answers = [fc.answer for fc in flashcards]
        else:
            # Extract questions and answers from tuples
            questions = [q for q, a in flashcards]
            answers = [a for q, a in flashcards]
        
        # Compute TF-IDF embeddings for questions
        question_embeddings = semantic_processor.vectorizer.fit_transform(questions).toarray()
        
        # Find similar questions
        similarity_matrix = cosine_similarity(question_embeddings)
        
        # Find duplicates
        duplicates = set()
        for i in range(len(questions)):
            for j in range(i + 1, len(questions)):
                if similarity_matrix[i][j] > similarity_threshold:
                    # Keep the one with longer answer (more detailed)
                    if len(answers[i]) < len(answers[j]):
                        duplicates.add(i)
                    else:
                        duplicates.add(j)
        
        # Remove duplicates
        deduplicated = [flashcards[i] for i in range(len(flashcards)) if i not in duplicates]
        
        print(f"🔄 Removed {len(flashcards) - len(deduplicated)} duplicate flashcards")
        return deduplicated
        
    except Exception as e:
        print(f"Warning: Could not remove duplicates: {e}")
        return flashcards

# =====================
# Test Deck Creation
# =====================

def create_test_apkg_deck(output_path='test_deck.apkg'):
    # Sample flashcards with image occlusion
    sample_flashcards = [
        {
            'type': 'image_occlusion',
            'question_img': 'apkg_media_check/occluded_slide10_img1.png',
            'answer_img': 'apkg_media_check/slide11_img1.jpg',
        },
        {
            'type': 'image_occlusion',
            'question_img': 'apkg_media_check/occluded_slide13_img1.png',
            'answer_img': 'apkg_media_check/slide12_img2.png',
        },
    ]
    
    # Export the sample flashcards to an APKG file
    export_flashcards_to_apkg(sample_flashcards, output_path=output_path)
    print(f"Test deck created at {output_path}")

# =====================
# Main Orchestration
# =====================

if __name__ == "__main__":
    # Directly create a test APKG deck
    create_test_apkg_deck()

# =====================
# Quality Control Classes
# =====================

@dataclass
class AudioMetadataForCard:
    audio_files: List[str]

@dc.dataclass if False else dataclass
class Flashcard:
    question: str
    answer: str
    level: int
    slide_number: int
    confidence: float = 0.0
    is_cloze: bool = False
    cloze_text: str = ""
    image_path: Optional[str] = None  # Path to last related image (if any)
    preview_text: str = ""  # Cleaned cloze text for UI preview
    raw_cloze_text: str = ""  # Original cloze text as received or generated
    card_type: str = ""  # e.g., 'occlusion' for masked image cards
    # Audio metadata for enhanced context
    audio_metadata: Optional[AudioMetadataForCard] = None
    emphasis_weight: float = 1.0  # Weight based on lecturer emphasis (0-2)
    time_allocation: float = 0.0  # Seconds spent on this concept

class QualityController:
    """Handles flashcard quality control and optimization"""
    
    def __init__(self):
        self.stop_words = set(stopwords.words('english'))
        self.vectorizer = TfidfVectorizer(
            max_features=1000,
            stop_words='english',
            ngram_range=(1, 2)
        )
        
    def extract_key_terms(self, text: str) -> List[str]:
        """Extract key medical terms from text"""
        # Remove common words and extract medical terminology
        words = word_tokenize(text.lower())
        key_terms = [word for word in words if word not in self.stop_words and len(word) > 2]
        return key_terms
    
    def calculate_similarity(self, text1: str, text2: str) -> float:
        """Calculate semantic similarity between two texts"""
        try:
            # Use TF-IDF for similarity calculation
            texts = [text1, text2]
            tfidf_matrix = self.vectorizer.fit_transform(texts)
            similarity = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:2])[0][0]
            return similarity
        except:
            # Fallback to sequence matcher
            return SequenceMatcher(None, text1.lower(), text2.lower()).ratio()
    
    def detect_repetition(self, cards: List[Flashcard], threshold: float = 0.7) -> List[Tuple[int, int]]:
        """Detect repetitive flashcards"""
        duplicates = []
        for i in range(len(cards)):
            for j in range(i + 1, len(cards)):
                # Check question similarity
                q_similarity = self.calculate_similarity(cards[i].question, cards[j].question)
                # Check answer similarity
                a_similarity = self.calculate_similarity(cards[i].answer, cards[j].answer)
                
                # If both question and answer are similar, mark as duplicate
                if q_similarity > threshold and a_similarity > threshold:
                    duplicates.append((i, j))
        
        return duplicates
    
    def is_too_wordy(self, text: str, max_sentences: int = 2, max_words: int = 25) -> bool:
        """Check if text is too wordy for a flashcard"""
        sentences = sent_tokenize(text)
        words = word_tokenize(text)
        
        return len(sentences) > max_sentences or len(words) > max_words
    
    def split_wordy_answer(self, question: str, answer: str) -> List[Tuple[str, str]]:
        """Split a wordy answer into multiple focused cards"""
        sentences = sent_tokenize(answer)
        if len(sentences) <= 2:
            return [(question, answer)]
        
        # Split into multiple cards
        cards = []
        for i, sentence in enumerate(sentences):
            if len(sentence.strip()) > 10:  # Only create cards for substantial sentences
                new_question = f"{question} (Part {i+1})"
                cards.append((new_question, sentence.strip()))
        
        return cards
    
    def is_shallow_card(self, question: str, answer: str, level: int) -> bool:
        """Detect if a card is too shallow for its level"""
        # Level 1 can be basic, but Level 2 should have reasoning
        if level == 1:
            # Check if it's just a definition without context
            definition_indicators = ['stands for', 'means', 'is defined as', 'refers to']
            has_context = any(indicator in question.lower() for indicator in ['why', 'how', 'what explains', 'what pattern'])
            return not has_context and any(indicator in question.lower() for indicator in definition_indicators)
        
        elif level == 2:
            # Level 2 should have reasoning words
            reasoning_indicators = ['why', 'how', 'explain', 'compare', 'interpret', 'pattern', 'suggests', 'indicates']
            return not any(indicator in question.lower() for indicator in reasoning_indicators)
        
        return False
    
    def enrich_shallow_card(self, question: str, answer: str, level: int) -> Tuple[str, str]:
        """Enrich a shallow card with more context"""
        if level == 1:
            # Add minimal context for Level 1
            if 'what does' in question.lower() and 'stand for' in question.lower():
                # Convert "What does X stand for?" to "What does X measure/represent?"
                question = question.replace('stand for', 'measure in clinical practice')
            elif 'what is' in question.lower() and len(answer.split()) < 5:
                # Add context for short definitions
                question = question.replace('What is', 'What is the clinical significance of')
        
        elif level == 2:
            # Add reasoning for Level 2
            if 'what is' in question.lower():
                question = question.replace('What is', 'What explains why')
            elif 'what does' in question.lower():
                question = question.replace('What does', 'What pattern does')
        
        return question, answer
    
    def detect_cloze_opportunities(self, question: str, answer: str) -> Tuple[bool, str]:
        """Detect if a card would be better as a cloze deletion"""
        # More aggressive patterns for medical content
        cloze_patterns = [
            r'(\d+(?:\.\d+)?%?)',  # Numbers and percentages
            r'([A-Z]{2,}(?:\d+)?)',  # Acronyms like FEV1, FVC, MCV, RDW
            r'(\w+ (?:and|or) \w+)',  # Lists
            r'(\w+ (?:is|are) \w+)',  # Definitions
            r'(\w+ (?:deficiency|disease|syndrome|anemia|hemolysis))',  # Medical conditions
            r'(\w+ (?:enzyme|protein|receptor|antibody|monoclonal))',  # Medical molecules
            r'(\w+ (?:count|level|concentration|volume|width))',  # Medical measurements
            r'(\w+ (?:infection|inflammation|damage|stress|response))',  # Medical processes
        ]
        
        # Also check for key medical terms that should be clozed
        key_medical_terms = [
            'MCV', 'RDW', 'MCH', 'MCHC', 'NRBC', 'CBC', 'PT', 'INR', 'aPTT',
            'hemoglobin', 'erythrocyte', 'neutrophil', 'reticulocyte', 'thrombocyte',
            'spherocytosis', 'pyruvate', 'dehydrogenase', 'haptoglobin', 'eculizumab'
        ]
        
        # Check patterns first
        for pattern in cloze_patterns:
            matches = re.findall(pattern, answer)
            if matches and len(matches) <= 4:  # Allow up to 4 blanks
                return True, self.create_cloze_text(answer, matches)
        
        # Check for key medical terms
        found_terms = []
        for term in key_medical_terms:
            if term.lower() in answer.lower():
                found_terms.append(term)
        
        if found_terms and len(found_terms) <= 3:
            return True, self.create_cloze_text(answer, found_terms)
        
        # If no specific patterns found, try to cloze the most important part
        # Look for the main concept (usually the first significant phrase)
        words = answer.split()
        if len(words) >= 4:
            # Find the first meaningful phrase (skip articles, prepositions)
            skip_words = {'the', 'a', 'an', 'is', 'are', 'of', 'in', 'to', 'for', 'with', 'by'}
            meaningful_words = [w for w in words if w.lower() not in skip_words and len(w) > 2]
            
            if meaningful_words:
                # Cloze the first meaningful concept
                first_concept = meaningful_words[0]
                if len(first_concept) > 3:  # Only cloze substantial words
                    return True, self.create_cloze_text(answer, [first_concept])
        
        return False, ""
    
    def create_cloze_text(self, answer: str, key_terms: List[str]) -> str:
        """Create cloze deletion text"""
        cloze_text = answer
        for i, term in enumerate(key_terms, 1):
            cloze_text = cloze_text.replace(term, f"{{{{c{i}::{term}}}}}", 1)
        return cloze_text
    
    def assess_depth_consistency(self, cards: List[Flashcard]) -> List[int]:
        """Assess if cards are at the right depth for their level"""
        inconsistent_cards = []
        
        for i, card in enumerate(cards):
            if card.level == 1 and self.is_too_deep_for_level1(card.question, card.answer):
                inconsistent_cards.append(i)
            elif card.level == 2 and self.is_too_shallow_for_level2(card.question, card.answer):
                inconsistent_cards.append(i)
        
        return inconsistent_cards
    
    def is_too_deep_for_level1(self, question: str, answer: str) -> bool:
        """Check if Level 1 card is too deep"""
        deep_indicators = ['pattern', 'interpret', 'compare', 'explain why', 'clinical reasoning', 'differential']
        return any(indicator in question.lower() for indicator in deep_indicators)
    
    def is_too_shallow_for_level2(self, question: str, answer: str) -> bool:
        """Check if Level 2 card is too shallow"""
        shallow_indicators = ['what does', 'stand for', 'define', 'name', 'list']
        return any(indicator in question.lower() for indicator in shallow_indicators)
    
    def fix_depth_inconsistency(self, card: Flashcard) -> Flashcard:
        """Fix depth inconsistency by adjusting level or content"""
        if card.level == 1 and self.is_too_deep_for_level1(card.question, card.answer):
            # Move to Level 2
            card.level = 2
        elif card.level == 2 and self.is_too_shallow_for_level2(card.question, card.answer):
            # Either move to Level 1 or enrich
            if len(card.answer.split()) < 10:
                card.level = 1
            else:
                # Enrich the question
                card.question, card.answer = self.enrich_shallow_card(card.question, card.answer, 2)
        
        return card

def is_image_relevant_for_occlusion(image_path: str, slide_text: str, api_key: str, model: str) -> bool:
    """
    Determine if an image is relevant for image occlusion flashcards.
    Filters out decorative, scene-setting, or non-medical images.
    
    Args:
        image_path: Path to the image file
        slide_text: Associated slide text for context
        api_key: OpenAI API key
        model: Model name to use for analysis
        
    Returns:
        True if the image contains relevant medical/clinical content for occlusion
"""
    try:
        print(f"[DEBUG] is_image_relevant_for_occlusion called for: {os.path.basename(image_path)}")
        
        # Check if the path is actually a file
        if not os.path.isfile(image_path):
            print(f"[WARN] Skipping image relevance check for {image_path}: not a file")
            return False
            
        # Read and encode the image
        with open(image_path, "rb") as image_file:
            ext = os.path.splitext(image_path)[1][1:]
            base64_image = base64.b64encode(image_file.read()).decode("utf-8")
        
        # Create the analysis prompt
        analysis_prompt = f"""
        Analyze this image and determine if it contains content suitable for medical flashcard creation via image occlusion.
        
        CONTEXT: This image is from a slide with the following text:
        "{slide_text[:500]}..."  # Truncated for brevity
        
        CRITICAL FILTERING RULES:
        REJECT the image if it contains:
        - Decorative elements (logos, decorative graphics, stock photos)
        - Scene-setting images (hospitals, doctors, patients in general)
        - Generic medical imagery (stethoscopes, medical symbols)
        - Navigation elements (arrows, buttons, icons)
        - Pure text slides (better handled as text flashcards)
        - Charts/graphs with no medical data
        - Generic illustrations without specific medical content
        
        ACCEPT the image if it contains:
        - Anatomical diagrams with labeled structures
        - Medical charts/graphs with clinical data
        - Drug mechanism diagrams
        - Pathological specimens or histological images
        - ECG traces, X-rays, or other medical imaging
        - Flowcharts of medical processes or decision trees
        - Tables with medical data, lab values, or drug information
        - Biochemical pathways or molecular diagrams
        - Clinical algorithms or protocols
        
        Respond with ONLY "RELEVANT" or "NOT_RELEVANT" followed by a brief explanation.
        """
        
        # Call the API for analysis
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": analysis_prompt},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/{ext};base64,{base64_image}"}
                        }
                    ]
                }
            ],
            max_tokens=100,
            temperature=0.1
        )
        
        result = response.choices[0].message.content.strip().upper()
        print(f"[DEBUG] Vision API response for {os.path.basename(image_path)}: {result}")
        
        # Parse the response
        if result.startswith("RELEVANT"):
            print(f"✅ Image {os.path.basename(image_path)} deemed relevant for occlusion")
            return True
        else:
            print(f"❌ Image {os.path.basename(image_path)} filtered out as not relevant")
            return False
            
    except Exception as e:
        print(f"⚠️ Error analyzing image relevance: {e}")
        # Default to False if analysis fails
        return False


def filter_relevant_images_for_occlusion(slide_images: List[List[str]], slide_texts: List[str], api_key: str, model: str) -> List[List[str]]:
    """
    Filter images to only include those relevant for image occlusion flashcards.
    
    Args:
        slide_images: List of image paths for each slide
        slide_texts: List of slide texts
        api_key: OpenAI API key
        model: Model name to use for analysis
        
    Returns:
        Filtered list of relevant images for each slide
    """
    print(f"[DEBUG] filter_relevant_images_for_occlusion called with {len(slide_images)} slides")
    
    if not slide_images or not slide_texts:
        print(f"[DEBUG] No slide images or texts provided: images={len(slide_images) if slide_images else 0}, texts={len(slide_texts) if slide_texts else 0}")
        return slide_images
    
    filtered_images = []
    
    for slide_idx, (images, slide_text) in enumerate(zip(slide_images, slide_texts)):
        print(f"[DEBUG] Processing slide {slide_idx + 1}: {len(images)} images, text length: {len(slide_text)}")
        relevant_images = []
        
        for image_path in images:
            print(f"[DEBUG] Checking image relevance: {os.path.basename(image_path)}")
            if is_image_relevant_for_occlusion(image_path, slide_text, api_key, model):
                relevant_images.append(image_path)
                print(f"[DEBUG] ✅ Image {os.path.basename(image_path)} deemed relevant")
            else:
                print(f"[DEBUG] ❌ Image {os.path.basename(image_path)} filtered out")
        
        filtered_images.append(relevant_images)
        
        if images and not relevant_images:
            print(f"📝 Slide {slide_idx + 1}: All {len(images)} images filtered out as not relevant")
        elif images:
            print(f"📝 Slide {slide_idx + 1}: {len(relevant_images)}/{len(images)} images deemed relevant")
    
    total_relevant = sum(len(images) for images in filtered_images)
    print(f"[DEBUG] Total relevant images found: {total_relevant}")
    return filtered_images

def export_occlusion_flashcards_to_csv(flashcard_entries, csv_path):
    """
    Export a list of occlusion flashcard entries to a CSV file for Anki import.
    Each entry should have 'question_image_base64' and 'answer_image_base64'.
    """
    with open(csv_path, 'w', newline='', encoding='utf-8') as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(["Front", "Back"])
        for entry in flashcard_entries:
            front = f"<img src='{os.path.basename(entry['question_image_path'])}'>"
            back  = f"<img src='{os.path.basename(entry['answer_image_path'])}'>"
            writer.writerow([front, back])

# Example usage of batch_generate_image_occlusion_flashcards
# image_paths = ['path/to/image1.png', 'path/to/image2.png']  # Replace with actual image paths
# export_dir = 'path/to/export_dir'  # Replace with actual export directory
# flashcard_entries = batch_generate_image_occlusion_flashcards(image_paths, export_dir)

# Process flashcard_entries as needed
# for entry in flashcard_entries:
#     print(f"Generated flashcard entry: {entry}")

# Use config values in batch_generate_image_occlusion_flashcards
# batch_generate_image_occlusion_flashcards(image_paths, export_dir, conf_threshold=config['conf_threshold'], mask_method='rectangle')

# batch_generate_image_occlusion_flashcards(image_paths, export_dir, conf_threshold=config['conf_threshold'], mask_method='rectangle') 