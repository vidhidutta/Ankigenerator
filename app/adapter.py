from __future__ import annotations

"""
Adapter for slide-to-cards conversion.

Exposes a single function:
  extract_cards_from_ppt(input_path: str) -> list[tuple[str, str]]

This calls the existing converter in `flashcard_generator.py`, then normalizes
results to a list of (question, answer) tuples.

Environment variables used by the underlying converter (detected from the codebase):
- OPENAI_API_KEY
- GOOGLE_APPLICATION_CREDENTIALS (used by image/vision providers; not required for basic text cards)

If the converter fails or returns nothing, this function returns an empty list
so the caller can provide a fallback card.
"""

from typing import Any, Iterable, List, Tuple
import os
import tempfile
import traceback
import json
import inspect
from pptx import Presentation
from openai import OpenAI


def _flatten(items: Iterable[Any]) -> list[Any]:
    flat: list[Any] = []
    for it in items or []:
        if isinstance(it, list):
            flat.extend(_flatten(it))
        else:
            flat.append(it)
    return flat


def _to_pairs(obj: Any) -> list[Tuple[str, str]]:
    pairs: list[Tuple[str, str]] = []

    # Already a list of pairs
    if isinstance(obj, list) and obj and isinstance(obj[0], tuple) and len(obj[0]) == 2:
        return [(str(q), str(a)) for (q, a) in obj]

    # DataFrame-like (duck typing; avoid hard pandas dependency)
    try:
        cols = getattr(obj, "columns", None)
        if cols is not None:
            col_names = [str(c).lower() for c in list(cols)]
            qcol = "q" if "q" in col_names else ("question" if "question" in col_names else None)
            acol = "a" if "a" in col_names else ("answer" if "answer" in col_names else None)
            if qcol and acol:
                # Iterate rows without importing pandas explicitly
                for idx in range(len(getattr(obj, "index", []))):
                    row = obj.iloc[idx]  # type: ignore[attr-defined]
                    q = str(row[qcol]) if qcol in row else ""
                    a = str(row[acol]) if acol in row else ""
                    if q.strip() and a.strip():
                        pairs.append((q.strip(), a.strip()))
                return pairs
    except Exception:
        pass

    # List of dicts
    if isinstance(obj, list) and obj and isinstance(obj[0], dict):
        for d in obj:
            q = d.get("q") if "q" in d else d.get("question")
            a = d.get("a") if "a" in d else d.get("answer")
            if isinstance(q, str) and isinstance(a, str) and q.strip() and a.strip():
                pairs.append((q.strip(), a.strip()))
        return pairs

    # List of objects with .question/.answer (e.g., Flashcard)
    if isinstance(obj, list) and obj and hasattr(obj[0], "question") and hasattr(obj[0], "answer"):
        for fc in obj:
            try:
                q = str(getattr(fc, "question", ""))
                a = str(getattr(fc, "answer", ""))
                if q.strip() and a.strip():
                    pairs.append((q.strip(), a.strip()))
            except Exception:
                continue
        return pairs

    # Unknown shape
    return []


def _pptx_text_blocks(path: str) -> list[str]:
    """Extract text blocks from PPTX using python-pptx (lightweight)."""
    texts = []
    prs = Presentation(path)
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if parts:
            texts.append("\n".join(parts))
    return texts


def _openai_cards_from_texts(blocks: list[str], max_cards: int = 30) -> list[tuple[str,str]]:
    """Generate flashcards from text blocks using OpenAI directly (lightweight)."""
    if not blocks:
        return []
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    prompt = (
        "You convert lecture notes into Anki Basic flashcards.\n"
        "Return ONLY valid JSON: a list of objects {\"q\":..., \"a\":...}.\n"
        f"Max {max_cards} cards. Use concise, clinically relevant Q/A. No markdown.\n"
    )
    joined = "\n\n--- SLIDES ---\n\n" + "\n\n---\n\n".join(blocks[:15])  # limit context
    r = client.chat.completions.create(
        model=os.getenv("OJAMED_OPENAI_MODEL","gpt-4o-mini"),
        messages=[{"role":"system","content":prompt},
                  {"role":"user","content":joined}],
        temperature=float(os.getenv("OJAMED_TEMPERATURE","0.2"))
    )
    text = r.choices[0].message.content.strip()
    try:
        data = json.loads(text)
    except Exception:
        # attempt to extract JSON substring
        import re
        m = re.search(r'\[.*\]', text, flags=re.S)
        data = json.loads(m.group(0)) if m else []
    cards = []
    if isinstance(data, list):
        for it in data:
            if isinstance(it, dict):
                q = it.get("q") or it.get("question")
                a = it.get("a") or it.get("answer")
                if q and a:
                    cards.append((str(q).strip(), str(a).strip()))
            elif isinstance(it, (list, tuple)) and len(it) >= 2:
                cards.append((str(it[0]).strip(), str(it[1]).strip()))
    return cards[:max_cards]


def extract_cards_from_ppt(input_path: str):
    """
    Extract flashcards from PowerPoint file using holistic medical analysis
    
    This enhanced version now provides:
    1. Traditional flashcards (Level 1 & 2)
    2. Comprehensive medical concept analysis
    3. Visual mind maps showing relationships
    4. Knowledge gap filling with expert explanations
    5. Clinical pearls and learning objectives
    """
    
    # DEMO BYPASS
    if os.getenv("OJAMED_FORCE_DEMO") == "1":
        print("[OjaMed][ADAPTER] DEMO MODE -> returning 3 static cards")
        return [
            ("What drug class is furosemide?", "Loop diuretic"),
            ("Main adverse effect?", "Hypokalemia"),
            ("Contraindicated with?", "Sulfa allergy (relative)"),
        ]
    
    try:
        # Check if holistic analysis is enabled
        if os.getenv("OJAMED_HOLISTIC_ANALYSIS", "1") == "1":
            print("[OjaMed][ADAPTER] 🧠 Holistic medical analysis enabled")
            
            # Import holistic analyzer
            try:
                from holistic_medical_analyzer import HolisticMedicalAnalyzer
                from medical_notes_pdf_generator import generate_medical_notes_pdf
                
                # Initialize holistic analyzer
                api_key = os.getenv("OPENAI_API_KEY")
                if not api_key:
                    print("[OjaMed][ADAPTER] ⚠️ No OpenAI API key found, falling back to basic extraction")
                    return _basic_extraction_fallback(input_path)
                
                analyzer = HolisticMedicalAnalyzer(api_key)
                
                # Extract text content from PPTX
                lecture_content = _extract_full_lecture_content(input_path)
                if not lecture_content:
                    print("[OjaMed][ADAPTER] ⚠️ No content extracted, falling back to basic extraction")
                    return _basic_extraction_fallback(input_path)
                
                print(f"[OjaMed][ADAPTER] 📚 Extracted {len(lecture_content)} characters of lecture content")
                
                # Perform holistic analysis
                lecture_title = os.path.basename(input_path).replace('.pptx', '').replace('.ppt', '')
                holistic_analysis = analyzer.analyze_lecture_holistically(lecture_content, lecture_title)
                
                print(f"[OjaMed][ADAPTER] 🧠 Holistic analysis complete:")
                print(f"  • {len(holistic_analysis.concepts)} concepts identified")
                print(f"  • {len(holistic_analysis.mind_maps)} mind maps generated")
                print(f"  • {len(holistic_analysis.filled_gaps)} knowledge gaps filled")
                print(f"  • {len(holistic_analysis.clinical_pearls)} clinical pearls extracted")
                
                # Generate comprehensive PDF notes
                try:
                    pdf_path = generate_medical_notes_pdf(holistic_analysis, f"{lecture_title}_comprehensive_notes.pdf")
                    print(f"[OjaMed][ADAPTER] 📄 Comprehensive notes PDF generated: {pdf_path}")
                    
                    # Store PDF path for pipeline to include in ZIP
                    os.environ["OJAMED_COMPREHENSIVE_NOTES_PDF"] = pdf_path
                    
                except Exception as e:
                    print(f"[OjaMed][ADAPTER] ⚠️ PDF generation failed: {e}")
                    traceback.print_exc()
                
                # Convert concepts to flashcards
                cards = _convert_concepts_to_flashcards(holistic_analysis)
                print(f"[OjaMed][ADAPTER] 🎯 Generated {len(cards)} flashcards from holistic analysis")
                
                return cards
                
            except ImportError as e:
                print(f"[OjaMed][ADAPTER] ⚠️ Holistic analysis modules not available: {e}")
                print("[OjaMed][ADAPTER] Falling back to basic extraction")
                return _basic_extraction_fallback(input_path)
        
        else:
            print("[OjaMed][ADAPTER] Holistic analysis disabled, using basic extraction")
            return _basic_extraction_fallback(input_path)
            
    except Exception as e:
        print("[OjaMed][ADAPTER] Error in holistic analysis:", repr(e))
        traceback.print_exc()
        
        # When debugging, let the API surface the error upstream
        if os.getenv("OJAMED_DEBUG") == "1":
            raise
        
        print("[OjaMed][ADAPTER] Falling back to basic extraction")
        return _basic_extraction_fallback(input_path)

def _extract_full_lecture_content(input_path: str) -> str:
    """Extract all text content from PowerPoint file"""
    try:
        prs = Presentation(input_path)
        content = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if slide_content:
                content.append(f"--- SLIDE {slide_num + 1} ---")
                content.extend(slide_content)
                content.append("")  # Empty line between slides
        
        return "\n".join(content)
        
    except Exception as e:
        print(f"[OjaMed][ADAPTER] Error extracting lecture content: {e}")
        return ""

def _convert_concepts_to_flashcards(holistic_analysis) -> list:
    """Convert holistic analysis concepts to flashcard format"""
    cards = []
    
    # Convert concepts to basic flashcards
    for concept in holistic_analysis.concepts:
        # Level 1 card (basic recall)
        if concept.definition:
            question = f"What is {concept.name}?"
            answer = concept.definition
            cards.append((question, answer))
        
        # Level 2 card (clinical application) if clinical relevance exists
        if concept.clinical_relevance:
            question = f"What is the clinical significance of {concept.name}?"
            answer = concept.clinical_relevance
            cards.append((question, answer))
        
        # Relationship cards
        for related in concept.relationships[:2]:  # Limit to 2 relationships per concept
            question = f"How does {concept.name} relate to {related}?"
            answer = f"{concept.name} and {related} are related concepts in this lecture"
            cards.append((question, answer))
    
    # Add clinical pearls as flashcards
    for pearl in holistic_analysis.clinical_pearls:
        question = "What is a key clinical pearl from this lecture?"
        answer = pearl
        cards.append((question, answer))
    
    # Add learning objectives as flashcards
    for objective in holistic_analysis.learning_objectives[:3]:  # Limit to 3
        question = "What is a learning objective for this lecture?"
        answer = objective
        cards.append((question, answer))
    
    return cards

def _basic_extraction_fallback(input_path: str):
    """Fallback to basic extraction when holistic analysis is not available"""
    
    # If light mode, skip importing flashcard_generator/sklearn entirely
    if os.getenv("OJAMED_DISABLE_SKLEARN") == "1":
        try:
            blocks = _pptx_text_blocks(input_path)
            cards = _openai_cards_from_texts(blocks, max_cards=int(os.getenv("OJAMED_MAX_CARDS","30")))
            print(f"[OjaMed][ADAPTER][LIGHT] -> {len(cards)} cards")
            return cards
        except Exception as e:
            print("[OjaMed][ADAPTER][LIGHT] failed:", repr(e))
            traceback.print_exc()
            return []
    
    # else use the existing heavy path (unchanged)
    try:
        from flashcard_generator import (
            extract_text_from_pptx,
            extract_images_from_pptx,
            generate_enhanced_flashcards_with_progress,
        )
    except Exception as e:
        print("[OjaMed][ADAPTER] import failed:", repr(e))
        traceback.print_exc()
        return []
    
    try:
        texts = extract_text_from_pptx(input_path)
        images = extract_images_from_pptx(input_path)
        cards_obj = generate_enhanced_flashcards_with_progress(
            texts, images,
            model=os.getenv("OJAMED_OPENAI_MODEL","gpt-4o-mini"),
            temperature=float(os.getenv("OJAMED_TEMPERATURE","0.2")),
            max_tokens=int(os.getenv("OJAMED_MAX_TOKENS","500")),
            api_key=os.getenv("OPENAI_API_KEY"),
        )
    except Exception as e:
        print("[OjaMed][ADAPTER] converter raised:", repr(e))
        traceback.print_exc()
        if os.getenv("OJAMED_DEBUG") == "1":
            raise
        return []
    
    # normalize as before (keep existing normalization code) and return list[(q,a)]
    cards = []
    try:
        if isinstance(cards_obj, dict) and "cards" in cards_obj:
            cards_obj = cards_obj["cards"]
        if hasattr(cards_obj, "to_dict"):
            cards_obj = cards_obj.to_dict(orient="records")
        if isinstance(cards_obj, list):
            for row in cards_obj:
                if isinstance(row, (tuple, list)) and len(row) >= 2:
                    q, a = row[0], row[1]
                elif isinstance(row, dict):
                    q = row.get("q") or row.get("question") or row.get("Question")
                    a = row.get("a") or row.get("answer") or row.get("Answer")
                else:
                    continue
                if q and a:
                    cards.append((str(q).strip(), str(a).strip()))
    except Exception as e:
        print("[OjaMed][ADAPTER] normalization failed:", repr(e))
        traceback.print_exc()
    
    return cards


