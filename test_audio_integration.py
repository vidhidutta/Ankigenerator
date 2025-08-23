#!/usr/bin/env python3
"""
Test script for audio integration feature
Demonstrates audio processing and enhanced flashcard generation
"""

import os
import tempfile
import json
import pytest
from typing import List, Dict

AUDIO_PATH = os.path.join('tests', 'data', 'sample_lecture_3min.mp3')
PPTX_PATH = os.path.join('tests', 'data', 'sample_4slides.pptx')

pytestmark = pytest.mark.audio_e2e

# ------------------------
# Self-contained fixtures
# ------------------------

def ensure_sample_audio(target_mp3: str) -> str:
    """Create a ~3-minute mono 16kHz audio file if missing.
    Prefer MP3 via pydub+ffmpeg; fallback to WAV via soundfile if needed.
    Returns the actual path created (may be .wav).
    """
    os.makedirs(os.path.dirname(target_mp3), exist_ok=True)
    if os.path.isfile(target_mp3):
        return target_mp3
    # Try pydub + ffmpeg
    actual = target_mp3
    try:
        from pydub import AudioSegment
        from pydub.generators import Sine
        import shutil as _sh
        if not _sh.which('ffmpeg'):
            raise RuntimeError('ffmpeg missing')
        sr = 16000
        total_ms = 180_000
        seg = AudioSegment.silent(duration=0)
        # build 6 x 30s tones with small silent gaps to mimic structure
        for i in range(6):
            tone = Sine(440 + 20*i).to_audio_segment(duration=28_000).set_frame_rate(sr).set_channels(1)
            gap = AudioSegment.silent(duration=2_000)
            seg += tone + gap
        # pad to 180s
        if len(seg) < total_ms:
            seg += AudioSegment.silent(duration=total_ms - len(seg))
        seg.export(target_mp3, format='mp3', bitrate='64k')
        return actual
    except Exception:
        # Fallback: write WAV via numpy + soundfile
        import numpy as np
        import soundfile as sf
        sr = 16000
        total_s = 180
        t = np.linspace(0, total_s, int(sr*total_s), endpoint=False)
        x = 0.1*np.sin(2*np.pi*440*t).astype('float32')
        wav_path = os.path.splitext(target_mp3)[0] + '.wav'
        sf.write(wav_path, x, sr)
        return wav_path

def ensure_sample_pptx(target_pptx: str) -> str:
    """Create a 4-slide PPTX with simple text bullets if missing."""
    os.makedirs(os.path.dirname(target_pptx), exist_ok=True)
    if os.path.isfile(target_pptx):
        return target_pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    titles = [
        "Intro to Pharmacology",
        "Beta Blockers",
        "Adverse Effects",
        "Clinical Cases",
    ]
    bullets = [
        ["Drug classes", "Mechanisms", "Clinical use"],
        ["Propranolol", "Beta-1 blockade", "Indications"],
        ["Bradycardia", "Fatigue", "Contraindications"],
        ["HTN case", "Angina case", "Decisions"],
    ]
    for i in range(4):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = titles[i]
        tx = slide.shapes.placeholders[1].text_frame
        tx.clear()
        for b in bullets[i]:
            p = tx.add_paragraph()
            p.text = b
            p.level = 0
    prs.save(target_pptx)
    return target_pptx

# ------------------------
# Tests
# ------------------------

@pytest.mark.audio_e2e
def test_e2e_audio_pipeline(tmp_path):
    from audio_processor import AudioProcessor
    from flashcard_generator import generate_enhanced_flashcards_with_progress, export_flashcards_to_apkg
    audio_path = ensure_sample_audio(AUDIO_PATH)
    pptx_path = ensure_sample_pptx(PPTX_PATH)

    # Extract slide texts
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide_texts = []
    for s in prs.slides:
        texts = []
        for shape in s.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
        slide_texts.append("\n".join(texts))

    ap = AudioProcessor(model_name="small")
    bundle = ap.build_audio_bundle(audio_path, slide_texts, diarization_enabled=False, alignment_mode="semantic+keyword")

    confidences = [w.confidence for w in bundle.slide_windows]
    high = sum(1 for c in confidences if c >= 0.5)
    assert high >= 3, f"Expected >=3 high-confidence mappings, got {high}"

    clips, _warnings = ap.create_audio_clips_for_bundle(bundle, output_dir=str(tmp_path), per_slide_max=2, min_clip_sec=6, max_clip_sec=6)
    total_clips = sum(len(v) for v in clips.values())
    assert total_clips >= 1, "Expected at least one audio clip"

    # Validate durations
    import librosa
    for files in clips.values():
        for f in files:
            d = float(librosa.get_duration(filename=f))
            assert 5.75 <= d <= 15.25, f"Clip duration out of bounds: {d} for {f}"

    flashcards, analysis = generate_enhanced_flashcards_with_progress(slide_texts, [], os.getenv('OPENAI_API_KEY','sk-test'), 'gpt-4o-mini', 512, 0.2, audio_bundle=bundle)
    out_apkg = tmp_path / 'golden_sample.apkg'
    export_flashcards_to_apkg(flashcards, str(out_apkg))
    assert out_apkg.exists() and out_apkg.stat().st_size > 0

    # Check referenced media filenames exist among created clips
    sound_refs = []
    for fc in flashcards:
        if hasattr(fc, 'audio_metadata') and fc.audio_metadata:
            for apath in fc.audio_metadata.audio_files:
                sound_refs.append(os.path.basename(apath))
    assert sound_refs, "No audio filenames referenced in cards"
    for fn in sound_refs:
        found = any(fn == os.path.basename(f) for lst in clips.values() for f in lst)
        assert found, f"Referenced sound file not present in created clips: {fn}"

@pytest.mark.audio_e2e
def test_e2e_no_audio_still_exports(tmp_path):
    from flashcard_generator import generate_enhanced_flashcards_with_progress, export_flashcards_to_apkg
    pptx_path = ensure_sample_pptx(PPTX_PATH)
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide_texts = []
    for s in prs.slides:
        texts = []
        for shape in s.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
        slide_texts.append("\n".join(texts))

    flashcards, analysis = generate_enhanced_flashcards_with_progress(slide_texts, [], os.getenv('OPENAI_API_KEY','sk-test'), 'gpt-4o-mini', 512, 0.2, audio_bundle=None)
    out_apkg = tmp_path / 'no_audio.apkg'
    export_flashcards_to_apkg(flashcards, str(out_apkg))
    assert out_apkg.exists() and out_apkg.stat().st_size > 0

@pytest.mark.audio_e2e
def test_e2e_too_short_window_warning(tmp_path):
    from audio_processor import AudioProcessor
    audio_path = ensure_sample_audio(AUDIO_PATH)
    pptx_path = ensure_sample_pptx(PPTX_PATH)
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide_texts = []
    for s in prs.slides:
        texts = []
        for shape in s.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
        slide_texts.append("\n".join(texts))
    ap = AudioProcessor(model_name="small")
    bundle = ap.build_audio_bundle(audio_path, slide_texts, diarization_enabled=False, alignment_mode="semantic+keyword")
    # Force clip bounds to be strict: min 12s so some windows will be skipped
    clips, warnings_map = ap.create_audio_clips_for_bundle(bundle, output_dir=str(tmp_path), per_slide_max=1, min_clip_sec=12, max_clip_sec=12)
    any_warn = any("Skipped:" in msg for msgs in warnings_map.values() for msg in msgs)
    assert any_warn, "Expected a 'Skipped:' warning for too-short window"

@pytest.mark.audio_e2e
def test_diarization_filters_non_lecturer(monkeypatch, tmp_path):
    from audio_processor import AudioProcessor
    from audio_types import TranscriptSegment, SlideAudioWindow, AudioBundle
    audio_path = ensure_sample_audio(AUDIO_PATH)
    segs = [
        TranscriptSegment(start=0.0, end=5.0, text="lecturer intro", speaker="LECTURER", emphasis=0.7),
        TranscriptSegment(start=5.0, end=10.0, text="student question", speaker="OTHER", emphasis=0.9),
        TranscriptSegment(start=10.0, end=15.0, text="lecturer answer", speaker="LECTURER", emphasis=0.8),
    ]
    win = SlideAudioWindow(slide_id=0, window=(0.0, 15.0), confidence=0.8, segments=segs)
    bundle = AudioBundle(audio_path=audio_path, segments=segs, slide_windows=[win], diarization_applied=True)
    ap = AudioProcessor(model_name="small")
    clips, warnings_map = ap.create_audio_clips_for_bundle(bundle, output_dir=str(tmp_path), per_slide_max=2, min_clip_sec=6, max_clip_sec=6)
    for files in clips.values():
        for f in files:
            import re
            m = re.search(r"_(\d+)-(\d+)\.mp3$", f)
            assert m, "Clip filename format unexpected"
            s_ms, e_ms = int(m.group(1)), int(m.group(2))
            s, e = s_ms/1000.0, e_ms/1000.0
            assert not (s >= 5.0 and e <= 10.0), "Non-LECTURER-only segment should be filtered"

@pytest.mark.audio_e2e
def test_diarization_unavailable_fallback(monkeypatch):
    from audio_processor import AudioProcessor
    audio_path = ensure_sample_audio(AUDIO_PATH)
    def _raise(*args, **kwargs):
        raise RuntimeError("pyannote unavailable")
    ap = AudioProcessor(model_name="small")
    monkeypatch.setattr(ap, '_diarize_speakers', _raise)
    from pptx import Presentation
    pptx_path = ensure_sample_pptx(PPTX_PATH)
    prs = Presentation(pptx_path)
    slide_texts = []
    for s in prs.slides:
        texts = []
        for shape in s.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
        slide_texts.append("\n".join(texts))
    bundle = ap.build_audio_bundle(audio_path, slide_texts, diarization_enabled=True, alignment_mode="semantic+keyword")
    assert bundle is not None
    assert getattr(bundle, 'diarization_applied', False) is False

@pytest.mark.audio_e2e
def test_emphasis_toggles_order(tmp_path):
    from audio_processor import AudioProcessor
    from flashcard_generator import generate_enhanced_flashcards_with_progress
    audio_path = ensure_sample_audio(AUDIO_PATH)
    pptx_path = ensure_sample_pptx(PPTX_PATH)
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide_texts = []
    for s in prs.slides:
        texts = []
        for shape in s.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
        slide_texts.append("\n".join(texts))
    ap = AudioProcessor(model_name="small")
    bundle = ap.build_audio_bundle(audio_path, slide_texts, diarization_enabled=False, alignment_mode="semantic+keyword")
    # Emphasis ON
    setattr(bundle, 'use_emphasis', True)
    cards_on, _ = generate_enhanced_flashcards_with_progress(slide_texts, [], os.getenv('OPENAI_API_KEY','sk-test'), 'gpt-4o-mini', 512, 0.2, audio_bundle=bundle)
    order_on = [getattr(c, 'slide_number', -1) for c in cards_on]
    # Emphasis OFF
    setattr(bundle, 'use_emphasis', False)
    cards_off, _ = generate_enhanced_flashcards_with_progress(slide_texts, [], os.getenv('OPENAI_API_KEY','sk-test'), 'gpt-4o-mini', 512, 0.2, audio_bundle=bundle)
    order_off = [getattr(c, 'slide_number', -1) for c in cards_off]
    assert order_on != order_off, "Expected card order to change when emphasis is toggled" 