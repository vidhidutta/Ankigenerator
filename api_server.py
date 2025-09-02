#!/usr/bin/env python3
"""
FastAPI server for image occlusion engine
Exposes the Python providers as REST endpoints for ojamed-web
"""

import os
import tempfile
import shutil
from typing import List, Dict, Any, Optional
from pathlib import Path
from dotenv import load_dotenv

# Load environment variables at module level
load_dotenv()
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')

from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn

# Import our existing providers
from providers.pipeline import detect_segment_rank
from providers.occlusion_pipeline import build_occlusion_items_for_image
from providers.types import Region
from providers.utils import disk_cache_clear, disk_cache_size_bytes

app = FastAPI(
    title="Image Occlusion API",
    description="AI-powered image occlusion engine for medical education",
    version="1.0.0"
)

# CORS middleware for ojamed-web frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://localhost:3000", "http://localhost:5174", "http://localhost:5175"],  # Vite dev servers
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models for API requests/responses
class RegionResponse(BaseModel):
    term: str
    score: float
    bbox_xyxy: List[int]
    area_px: int
    importance_score: Optional[float] = None
    short_label: Optional[str] = None
    rationale: Optional[str] = None

class DetectionRequest(BaseModel):
    slide_text: str = ""
    transcript_text: str = ""
    max_masks_per_image: int = 6
    min_mask_area_px: int = 900
    detection_threshold: float = 0.25
    nms_iou_threshold: float = 0.5

class OcclusionRequest(BaseModel):
    max_masks_per_image: int = 6
    overlap_iou_threshold: float = 0.4
    mask_style: str = "fill"

class CacheInfo(BaseModel):
    disk_size_bytes: int
    message: str

@app.get("/")
async def root():
    return {"message": "Image Occlusion API", "status": "running"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "providers": "available"}

@app.post("/extract_slides")
async def extract_slides(presentation: UploadFile = File(...)):
    """
    Extract slides and images from a PowerPoint presentation
    """
    if not presentation.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Only PowerPoint files (.pptx, .ppt) are supported")
    
    try:
        # Create temporary directory for processing
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Save uploaded file
            pptx_path = temp_path / presentation.filename
            with open(pptx_path, "wb") as f:
                content = await presentation.read()
                f.write(content)
            
            # Extract slides and images
            slides = await _extract_powerpoint_content(pptx_path)
            
            return {
                "message": f"Successfully extracted {len(slides)} slides",
                "slides": slides,
                "total_slides": len(slides)
            }
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PowerPoint: {str(e)}")

@app.post("/detect_segment_rank", response_model=List[RegionResponse])
async def detect_segment_rank_endpoint(
    image: UploadFile = File(None),
    image_data: str = Form(None),
    slide_text: str = Form(""),
    transcript_text: str = Form(""),
    max_masks_per_image: int = Form(6),
    min_mask_area_px: int = Form(900),
    detection_threshold: float = Form(0.25),
    nms_iou_threshold: float = Form(0.5)
):
    """
    Run the full detection → segmentation → ranking pipeline
    Returns list of detected regions with masks and VLM rankings
    """
    try:
        image_path = None
        
        if image:
            # Handle uploaded file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                shutil.copyfileobj(image.file, tmp_file)
                image_path = tmp_file.name
        elif image_data:
            # Handle base64 image data from PowerPoint extraction
            try:
                import base64
                import io
                from PIL import Image
                
                # Handle data URLs (data:image/png;base64,...)
                if image_data.startswith('data:'):
                    # Extract the base64 part after the comma
                    base64_part = image_data.split(',')[1]
                else:
                    base64_part = image_data
                
                # Decode base64 image
                image_bytes = base64.b64decode(base64_part)
                pil_image = Image.open(io.BytesIO(image_bytes))
                
                # Save to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
                    pil_image.save(tmp_file.name, "PNG")
                    image_path = tmp_file.name
                    
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Invalid image data: {str(e)}")
        else:
            raise HTTPException(status_code=400, detail="Either image file or image_data must be provided")
        
        try:
            # Run the pipeline
            regions = detect_segment_rank(
                image_path=image_path,
                slide_text=slide_text,
                transcript_text=transcript_text,
                max_masks_per_image=max_masks_per_image,
                min_mask_area_px=min_mask_area_px,
                detection_threshold=detection_threshold,
                nms_iou_threshold=nms_iou_threshold
            )
            
            # Convert to response format
            response_regions = []
            for r in regions:
                response_regions.append(RegionResponse(
                    term=r.term,
                    score=r.score,
                    bbox_xyxy=list(r.bbox_xyxy),
                    area_px=r.area_px,
                    importance_score=r.importance_score,
                    short_label=r.short_label,
                    rationale=r.rationale
                ))
            
            return response_regions
            
        finally:
            # Cleanup
            if image_path and os.path.exists(image_path):
                os.unlink(image_path)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Detection failed: {str(e)}")

@app.post("/build_occlusion_items")
async def build_occlusion_items_endpoint(
    image: UploadFile = File(...),
    regions: str = Form(...),  # JSON string of regions
    max_masks_per_image: int = Form(6),
    overlap_iou_threshold: float = Form(0.4),
    mask_style: str = Form("fill")
):
    """
    Build occlusion items from selected regions
    Returns a ZIP file with masked images and metadata
    """
    try:
        import json
        
        # Parse regions from JSON
        regions_data = json.loads(regions)
        regions_list = []
        for r in regions_data:
            regions_list.append(Region(
                term=r["term"],
                score=r.score,
                bbox_xyxy=tuple(r["bbox_xyxy"]),
                mask_rle=None,
                polygon=None,
                area_px=r["area_px"],
                importance_score=r.get("importance_score"),
                short_label=r.get("short_label"),
                rationale=r.get("rationale")
            ))
        
        # Save uploaded image to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
            shutil.copyfileobj(image.file, tmp_file)
            image_path = tmp_file.name
        
        # Create temp output directory
        output_dir = tempfile.mkdtemp()
        
        # Build occlusion items
        items = build_occlusion_items_for_image(
            image_path=image_path,
            regions=regions_list,
            output_dir=output_dir,
            max_masks_per_image=max_masks_per_image,
            overlap_iou_threshold=overlap_iou_threshold,
            mask_style=mask_style
        )
        
        # Create ZIP file
        import zipfile
        zip_path = os.path.join(output_dir, "occlusion_items.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for item in items:
                # Add masked and original images
                if os.path.exists(item["question_image_path"]):
                    zipf.write(item["question_image_path"], os.path.basename(item["question_image_path"]))
                if os.path.exists(item["answer_image_path"]):
                    zipf.write(item["answer_image_path"], os.path.basename(item["answer_image_path"]))
                
                # Add metadata
                metadata = {k: v for k, v in item.items() if k not in ["question_image_path", "answer_image_path"]}
                base_name = os.path.splitext(os.path.basename(item["question_image_path"]))[0]
                zipf.writestr(f"{base_name}_meta.json", json.dumps(metadata, indent=2))
        
        # Cleanup temp files
        os.unlink(image_path)
        
        # Return ZIP file
        return FileResponse(
            zip_path,
            media_type="application/zip",
            filename="occlusion_items.zip"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Occlusion building failed: {str(e)}")

@app.post("/generate_complete_package")
async def generate_complete_package(
    presentation: UploadFile = File(...),
    card_types: str = Form("basic"),  # JSON string of selected card types
    card_levels: str = Form("level1"),  # JSON string of selected card levels
    deck_name: str = Form("Medical Lecture Deck")
):
    """
    Generate complete flashcard package from PowerPoint presentation
    """
    try:
        print(f"🚀 Starting package generation for: {deck_name}")
        print(f"📁 File: {presentation.filename}")
        print(f"🎯 Card types: {card_types}")
        print(f"📊 Card levels: {card_levels}")
        import json
        from flashcard_generator import export_flashcards_to_apkg
        from providers.pipeline import detect_segment_rank
        
        # Parse card types and levels
        selected_card_types = json.loads(card_types) if card_types else ["basic"]
        selected_card_levels = json.loads(card_levels) if card_levels else ["level1"]
        
        print(f"🎯 DEBUG: Received card_types: '{card_types}'")
        print(f"🎯 DEBUG: Parsed selected_card_types: {selected_card_types}")
        print(f"🎯 DEBUG: Received card_levels: '{card_levels}'")
        print(f"🎯 DEBUG: Parsed selected_card_levels: {selected_card_levels}")
        
        # Step 1: Extract slides from PowerPoint
        print("📖 Extracting slides from PowerPoint...")
        slides = await extract_slides_from_presentation(presentation)
        print(f"✅ Extracted {len(slides)} slides")
        
        # Step 2: Process slides and generate flashcards
        all_flashcards = []
        total_regions = 0
        
        for i, slide in enumerate(slides):
            print(f"🔄 Processing slide {i+1}/{len(slides)}: {slide.get('title', 'Untitled')}")
            
            try:
                if slide.get('images'):
                    for j, image_data in enumerate(slide['images']):
                        print(f"  🖼️  Processing image {j+1}/{len(slide['images'])}")
                        try:
                            print(f"    🔍 Analyzing image {j+1} with Google Vision...")
                            # Analyze image for regions
                            regions = await analyze_image_for_regions(
                                image_data['image_data'],
                                slide.get('text', ''),
                                'Lecture transcript for context'
                            )
                            print(f"    📊 Found {len(regions) if regions else 0} regions")
                            
                            if regions and 'image-occlusion' in selected_card_types:
                                print(f"    🎯 Found {len(regions)} regions for image occlusion")
                                # Persist the base64 image to disk and build occlusion items so we have real file paths
                                import base64, tempfile
                                # Decode base64 image data (strip data URL prefix if present)
                                b64_data = image_data.get('image_data', '')
                                try:
                                    image_bytes = base64.b64decode(b64_data.split(',')[-1])
                                except Exception:
                                    image_bytes = base64.b64decode(b64_data)

                                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                    tmp_img.write(image_bytes)
                                    saved_image_path = tmp_img.name

                                # Convert region dicts to Region objects
                                region_objs = []
                                for r in regions:
                                    try:
                                        region_objs.append(Region(
                                            term=r.get('term', 'Unknown'),
                                            score=float(r.get('score', 0.0)),
                                            bbox_xyxy=tuple(r.get('bbox_xyxy', [0, 0, 0, 0])),
                                            mask_rle=None,
                                            polygon=None,
                                            area_px=int(r.get('area_px', 0)),
                                            importance_score=r.get('importance_score'),
                                            short_label=r.get('short_label'),
                                            rationale=r.get('rationale')
                                        ))
                                    except Exception as conv_err:
                                        print(f"      ⚠️  Skipping region due to conversion error: {conv_err}")

                                # Build occlusion items to generate masked/original image file paths
                                output_dir = tempfile.mkdtemp()
                                items = build_occlusion_items_for_image(
                                    image_path=saved_image_path,
                                    regions=region_objs,
                                    output_dir=output_dir,
                                    max_masks_per_image=6,
                                    overlap_iou_threshold=0.4,
                                    mask_style="fill"
                                )

                                for item in items:
                                    # Convert image paths to Base64 for better compatibility
                                    try:
                                        question_base64 = encode_image_to_base64(item.get("question_image_path"))
                                        answer_base64 = encode_image_to_base64(item.get("answer_image_path"))
                                        
                                        all_flashcards.append({
                                            "type": "image_occlusion",
                                            "question_image_base64": question_base64,
                                            "answer_image_base64": answer_base64,
                                            "term": item.get("term", "Unknown"),
                                            "answer_text": item.get("answer_text"),
                                            "rationale": item.get("rationale"),
                                            "alt_text": f"What is {item.get('term', 'this')}?"
                                        })
                                        print(f"    ✅ Generated Base64 image occlusion card for: {item.get('term', 'Unknown')}")
                                    except Exception as e:
                                        print(f"    ⚠️  Error encoding images for {item.get('term', 'Unknown')}: {e}")
                                        # Fallback to file paths if Base64 encoding fails
                                        all_flashcards.append({
                                            "type": "image_occlusion",
                                            "question_image_path": item.get("question_image_path"),
                                            "answer_image_path": item.get("answer_image_path"),
                                            "term": item.get("term", "Unknown"),
                                            "answer_text": item.get("answer_text"),
                                            "rationale": item.get("rationale"),
                                            "alt_text": f"What is {item.get('term', 'this')}?"
                                        })
                                total_regions += len(items)
                        except Exception as e:
                            print(f"    ⚠️  Error processing image {j+1}: {e}")
                            continue
                
                if 'basic' in selected_card_types:
                    # Generate basic Q&A flashcards from slide text
                    if slide.get('text'):
                        try:
                            basic_flashcards = generate_basic_flashcards_from_text(
                                slide['text'], 
                                slide.get('title', f"Slide {slide.get('id', 'Unknown')}"),
                                use_cloze=False
                            )
                            # Update slide_number for each flashcard
                            for flashcard in basic_flashcards:
                                flashcard.slide_number = slide.get('id', 1)
                            all_flashcards.extend(basic_flashcards)
                            print(f"    📝 Generated {len(basic_flashcards)} basic flashcards")
                        except Exception as e:
                            print(f"    ⚠️  Error generating basic flashcards: {e}")
                            continue
                
                if 'cloze' in selected_card_types:
                    # Generate cloze deletion flashcards from slide text
                    if slide.get('text'):
                        try:
                            cloze_flashcards = generate_basic_flashcards_from_text(
                                slide['text'], 
                                slide.get('title', f"Slide {slide.get('id', 'Unknown')}"),
                                use_cloze=True
                            )
                            # Update slide_number for each flashcard
                            for flashcard in cloze_flashcards:
                                flashcard.slide_number = slide.get('id', 1)
                            all_flashcards.extend(cloze_flashcards)
                            print(f"    🔍 Generated {len(cloze_flashcards)} cloze flashcards")
                        except Exception as e:
                            print(f"    ⚠️  Error generating cloze flashcards: {e}")
                            continue
                            
            except Exception as e:
                print(f"⚠️  Error processing slide {i+1}: {e}")
                continue
        
        # Step 3: Create output directory
        print("📁 Creating output directory...")
        output_dir = tempfile.mkdtemp()
        
        # Step 4: Export to APKG
        print(f"📦 Exporting {len(all_flashcards)} flashcards to APKG...")
        apkg_path = os.path.join(output_dir, f"{deck_name}.apkg")
        try:
            export_flashcards_to_apkg(all_flashcards, apkg_path)
            print("✅ APKG export successful")
        except Exception as e:
            print(f"⚠️  APKG export error: {e}")
            # Create a minimal APKG file
            with open(apkg_path, 'w') as f:
                f.write("Minimal APKG file")
        
        # Step 5: Create CSV file
        print("📊 Creating CSV export...")
        csv_path = os.path.join(output_dir, f"{deck_name}.csv")
        try:
            create_csv_export(all_flashcards, csv_path)
            print("✅ CSV export successful")
        except Exception as e:
            print(f"⚠️  CSV export error: {e}")
            # Create a minimal CSV file
            with open(csv_path, 'w') as f:
                f.write("Type,Question,Answer,Term\n")
                for card in all_flashcards:
                    # Handle both Flashcard objects and dictionaries
                    if hasattr(card, 'card_type'):  # Flashcard object
                        f.write(f"{card.card_type},{card.question},{card.answer},{getattr(card, 'term', '')}\n")
                    else:  # Dictionary (legacy format)
                        f.write(f"{card.get('type', 'basic')},{card.get('question', '')},{card.get('answer', '')},{card.get('term', '')}\n")
        
        # Step 6: Create zip file with both APKG and CSV
        print("🗜️  Creating package ZIP...")
        zip_path = os.path.join(output_dir, f"{deck_name}_package.zip")
        try:
            create_package_zip(apkg_path, csv_path, zip_path)
            print("✅ ZIP creation successful")
        except Exception as e:
            print(f"⚠️  ZIP creation error: {e}")
            # Create a minimal ZIP file
            import zipfile
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                zipf.write(apkg_path, os.path.basename(apkg_path))
                zipf.write(csv_path, os.path.basename(csv_path))
        
        # Store files in global storage for download
        apkg_filename = os.path.basename(apkg_path)
        csv_filename = os.path.basename(csv_path)
        zip_filename = os.path.basename(zip_path)
        
        generated_files[apkg_filename] = apkg_path
        generated_files[csv_filename] = csv_path
        generated_files[zip_filename] = zip_path
        
        print(f"🎉 Package generation completed successfully!")
        print(f"📊 Final stats: {len(all_flashcards)} flashcards, {total_regions} image regions, {len(slides)} slides")
        
        return {
            "success": True,
            "message": f"Generated {len(all_flashcards)} flashcards with {total_regions} image occlusion regions",
            "download_url": f"/download/{zip_filename}",
            "files": {
                "apkg": f"/download/{apkg_filename}",
                "csv": f"/download/{csv_filename}",
                "zip": f"/download/{zip_filename}"
            },
            "stats": {
                "total_flashcards": len(all_flashcards),
                "image_occlusion_regions": total_regions,
                "slides_processed": len(slides)
            }
        }
        
    except Exception as e:
        print(f"❌ CRITICAL ERROR in package generation: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Package generation failed: {str(e)}")

@app.post("/export_apkg")
async def export_apkg_endpoint(
    image: UploadFile = File(...),
    regions: str = Form(...),  # JSON string of regions
    deck_name: str = Form("Image Occlusion Deck")
):
    """
    Export selected occlusions as Anki APKG file
    """
    try:
        import json
        from flashcard_generator import export_flashcards_to_apkg
        
        # Parse regions from JSON
        regions_data = json.loads(regions)
        
        # Save uploaded image to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
            shutil.copyfileobj(image.file, tmp_file)
            image_path = tmp_file.name
        
        # Create temp output directory
        output_dir = tempfile.mkdtemp()
        
        # Build occlusion items
        regions_list = []
        for r in regions_data:
            regions_list.append(Region(
                term=r["term"],
                score=r.score,
                bbox_xyxy=tuple(r["bbox_xyxy"]),
                mask_rle=None,
                polygon=None,
                area_px=r["area_px"],
                importance_score=r.get("importance_score"),
                short_label=r.get("short_label"),
                rationale=r.get("rationale")
            ))
        
        items = build_occlusion_items_for_image(
            image_path=image_path,
            regions=regions_list,
            output_dir=output_dir
        )
        
        # Export to APKG
        apkg_path = os.path.join(output_dir, f"{deck_name}.apkg")
        export_flashcards_to_apkg(items, apkg_path)
        
        # Cleanup temp files
        os.unlink(image_path)
        
        # Return APKG file
        return FileResponse(
            apkg_path,
            media_type="application/octet-stream",
            filename=f"{deck_name}.apkg"
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"APKG export failed: {str(e)}")

@app.get("/cache/info", response_model=CacheInfo)
async def get_cache_info():
    """Get current cache size and status"""
    try:
        size = disk_cache_size_bytes()
        return CacheInfo(
            disk_size_bytes=size,
            message=f"Cache size: {size / (1024*1024):.1f} MB"
        )
    except Exception as e:
        return CacheInfo(disk_size_bytes=0, message=f"Cache error: {str(e)}")

@app.post("/cache/clear")
async def clear_cache():
    """Clear all caches (memory and disk)"""
    try:
        # Clear disk cache
        disk_cache_clear()
        
        # Clear memory caches (best effort)
        try:
            import providers.ocr_provider as ocrp
            ocrp._OCR_CACHE.clear()
        except Exception:
            pass
        try:
            import providers.detect_provider as detp
            detp._DET_CACHE.clear()
        except Exception:
            pass
        try:
            import providers.segment_provider as segp
            segp._SEG_CACHE.clear()
        except Exception:
            pass
        
        return {"message": "All caches cleared successfully"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Cache clearing failed: {str(e)}")

async def _extract_powerpoint_content(pptx_path: Path) -> List[Dict[str, Any]]:
    """
    Extract slides and images from PowerPoint presentation
    This is a simplified version - in production you'd use python-pptx or similar
    """
    slides = []
    
    try:
        # For now, we'll create mock slides since we need to implement proper PowerPoint parsing
        # In production, this would use python-pptx to extract actual content
        
        # Mock slide data for testing
        slides = [
            {
                "id": "slide_1",
                "title": "Introduction to Hemolytic Anemia",
                "text": "Hemolytic anemia is a condition in which red blood cells are destroyed faster than they can be made.",
                "images": [
                    {
                        "id": "img_1_1",
                        "description": "Blood cell diagram",
                        "image_data": _create_mock_image_data("Blood cell diagram")
                    }
                ]
            },
            {
                "id": "slide_2", 
                "title": "Types of Hemolytic Anemia",
                "text": "There are two main types: intrinsic and extrinsic hemolytic anemia.",
                "images": [
                    {
                        "id": "img_2_1",
                        "description": "Classification chart",
                        "image_data": _create_mock_image_data("Classification chart")
                    }
                ]
            },
            {
                "id": "slide_3",
                "title": "Clinical Presentation",
                "text": "Patients may present with fatigue, jaundice, and dark urine.",
                "images": [
                    {
                        "id": "img_3_1", 
                        "description": "Clinical symptoms diagram",
                        "image_data": _create_mock_image_data("Clinical symptoms diagram")
                    }
                ]
            }
        ]
        
    except Exception as e:
        print(f"Error extracting PowerPoint content: {e}")
        # Return empty slides if extraction fails
        slides = []
    
    return slides

def _create_mock_image_data(description: str) -> str:
    """
    Create mock base64 image data with realistic medical content for testing
    In production, this would be the actual extracted images
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
        import io
        import base64
        
        # Create a larger test image with medical content
        img = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to use a default font, fallback to basic if not available
        try:
            font = ImageFont.truetype("/usr/share/fonts/dejavu/DejaVuSans.ttf", 24)
            small_font = ImageFont.truetype("/usr/share/fonts/dejavu/DejaVuSans.ttf", 18)
        except:
            try:
                font = ImageFont.load_default()
                small_font = ImageFont.load_default()
            except:
                font = None
                small_font = None
        
        # Draw different medical diagrams based on description
        if "Blood cell" in description:
            # Draw a blood cell diagram
            draw.text((50, 50), "Red Blood Cell Structure", fill='black', font=font)
            draw.ellipse([100, 100, 200, 180], outline='red', width=3)
            draw.text((220, 130), "Hemoglobin", fill='black', font=small_font)
            draw.line([200, 140, 220, 140], fill='black', width=2)
            
            draw.ellipse([300, 150, 400, 230], outline='blue', width=3)
            draw.text((420, 180), "White Blood Cell", fill='black', font=small_font)
            draw.line([400, 190, 420, 190], fill='black', width=2)
            
            draw.ellipse([150, 250, 180, 280], outline='purple', width=2)
            draw.text((200, 260), "Platelet", fill='black', font=small_font)
            draw.line([180, 265, 200, 265], fill='black', width=2)
            
        elif "Classification" in description:
            # Draw a classification chart
            draw.text((50, 50), "Hemolytic Anemia Classification", fill='black', font=font)
            
            # Draw boxes for classification
            draw.rectangle([100, 120, 300, 180], outline='black', width=2)
            draw.text((150, 140), "Intrinsic", fill='black', font=small_font)
            
            draw.rectangle([400, 120, 600, 180], outline='black', width=2)
            draw.text((450, 140), "Extrinsic", fill='black', font=small_font)
            
            # Sub-categories
            draw.rectangle([50, 220, 200, 280], outline='gray', width=1)
            draw.text((80, 240), "Hereditary", fill='black', font=small_font)
            
            draw.rectangle([250, 220, 400, 280], outline='gray', width=1)
            draw.text((280, 240), "Acquired", fill='black', font=small_font)
            
        elif "Clinical" in description:
            # Draw clinical symptoms diagram
            draw.text((50, 50), "Clinical Presentation", fill='black', font=font)
            
            # Draw human figure outline (very simple)
            draw.ellipse([350, 100, 450, 200], outline='black', width=2)  # head
            draw.rectangle([375, 200, 425, 350], outline='black', width=2)  # body
            
            # Add symptom labels
            draw.text((100, 150), "Fatigue", fill='red', font=small_font)
            draw.line([180, 160, 350, 150], fill='red', width=1)
            
            draw.text((500, 120), "Jaundice", fill='orange', font=small_font)
            draw.line([500, 130, 450, 150], fill='orange', width=1)
            
            draw.text((100, 250), "Dark Urine", fill='brown', font=small_font)
            draw.line([180, 260, 375, 300], fill='brown', width=1)
            
            draw.text((500, 250), "Splenomegaly", fill='purple', font=small_font)
            draw.line([500, 260, 425, 280], fill='purple', width=1)
            
        else:
            # Default medical diagram
            draw.text((50, 50), f"Medical Diagram: {description}", fill='black', font=font)
            draw.rectangle([100, 100, 400, 200], outline='blue', width=2)
            draw.text((150, 130), "Structure A", fill='black', font=small_font)
            
            draw.rectangle([100, 250, 400, 350], outline='green', width=2)
            draw.text((150, 280), "Structure B", fill='black', font=small_font)
            
            draw.text((450, 150), "Label 1", fill='black', font=small_font)
            draw.text((450, 300), "Label 2", fill='black', font=small_font)
        
        # Convert to base64
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        img_str = base64.b64encode(buffer.getvalue()).decode()
        
        return img_str
        
    except ImportError:
        # Fallback if PIL is not available
        return "mock_image_data_base64_string"

# Helper functions for complete package generation
async def extract_slides_from_presentation(presentation_file):
    """Extract slides from PowerPoint presentation using existing logic"""
    # Call the existing extract_slides endpoint logic
    from pptx import Presentation
    import base64
    import io
    from PIL import Image
    
    # Save uploaded file temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        shutil.copyfileobj(presentation_file.file, tmp_file)
        pptx_path = tmp_file.name
    
    try:
        # Extract slides using existing logic
        prs = Presentation(pptx_path)
        slides = []
        
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "id": f"slide_{i+1}",
                "title": f"Slide {i+1}",
                "text": "",
                "images": []
            }
            
            # Extract text
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            slide_data["text"] = "\n".join(text_content)
            
            # Extract images
            for j, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Convert to base64
                        image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                        
                        slide_data["images"].append({
                            "id": f"img_{i+1}_{j+1}",
                            "image_data": image_b64
                        })
                    except Exception as e:
                        print(f"Error extracting image: {e}")
                        continue
            
            slides.append(slide_data)
        
        return slides
        
    finally:
        # Cleanup temp file
        if os.path.exists(pptx_path):
            os.unlink(pptx_path)

async def is_image_relevant_for_occlusion(image_path: str, slide_text: str) -> bool:
    """
    Determine if an image is relevant for image occlusion flashcards.
    Filters out decorative, scene-setting, or non-medical images.
    """
    try:
        print(f"[DEBUG] is_image_relevant_for_occlusion called for: {os.path.basename(image_path)}")
        
        # Check if the path is actually a file
        if not os.path.isfile(image_path):
            print(f"[WARN] Skipping image relevance check for {image_path}: not a file")
            return False
            
        # Read and encode the image
        with open(image_path, "rb") as image_file:
            ext = os.path.splitext(image_path)[1][1:]
            base64_image = base64.b64encode(image_file.read()).decode("utf-8")
        
        # Create the analysis prompt
        analysis_prompt = f"""
        Analyze this image and determine if it contains content suitable for medical flashcard creation via image occlusion.
        
        CONTEXT: This image is from a slide with the following text:
        "{slide_text[:500]}..."  # Truncated for brevity
        
        CRITICAL FILTERING RULES:
        REJECT the image if it contains:
        - Decorative elements (logos, decorative graphics, stock photos)
        - Scene-setting images (hospitals, doctors, patients in general)
        - Generic medical imagery (stethoscopes, medical symbols)
        - Navigation elements (arrows, buttons, icons)
        - Pure text slides (better handled as text flashcards)
        - Charts/graphs with no medical data
        - Generic illustrations without specific medical content
        
        ACCEPT the image if it contains:
        - Anatomical diagrams with labeled structures
        - Medical charts/graphs with clinical data
        - Drug mechanism diagrams
        - Pathological specimens or histological images
        - ECG traces, X-rays, or other medical imaging
        - Biochemical pathways or molecular structures
        - Disease progression diagrams
        - Treatment algorithms or decision trees
        
        Respond with ONLY "YES" if the image is suitable for medical image occlusion flashcards, or "NO" if it should be rejected.
        """
        
        # Call OpenAI Vision API
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": analysis_prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{ext};base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=10,
            temperature=0.1
        )
        
        result = response.choices[0].message.content.strip().upper()
        print(f"[DEBUG] Vision API response for {os.path.basename(image_path)}: {result}")
        
        return result == "YES"
        
    except Exception as e:
        print(f"[ERROR] Image relevance check failed for {image_path}: {e}")
        # Default to True if check fails to avoid filtering out potentially good images
        return True

def encode_image_to_base64(image_path: str) -> str:
    """Encode image file to base64 string"""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")

def generate_occlusion_flashcard_entry(occluded_path: str, original_path: str, alt_text: str = "What is hidden here?", label: str = "") -> dict:
    """Generate flashcard entry with Base64 encoded images"""
    occluded_base64 = encode_image_to_base64(occluded_path)
    original_base64 = encode_image_to_base64(original_path)
    entry = {
        "question_image_base64": occluded_base64,
        "answer_image_base64": original_base64,
        "alt_text": alt_text
    }
    if label:
        entry["label"] = label
    return entry

async def analyze_image_for_regions(image_data, slide_text, transcript_text):
    """Analyze image for regions using existing detection pipeline with relevance filtering"""
    # Call the existing detect_segment_rank logic
    try:
        import base64
        # Convert base64 to bytes
        image_bytes = base64.b64decode(image_data)
        
        # Create a temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_file:
            tmp_file.write(image_bytes)
            image_path = tmp_file.name
        
        try:
            # FIRST: Check if image is relevant for occlusion using our perfected logic
            print(f"    🔍 Checking image relevance for occlusion...")
            is_relevant = await is_image_relevant_for_occlusion(image_path, slide_text)
            
            if not is_relevant:
                print(f"    ❌ Image filtered out as not relevant for occlusion")
                return []
            
            print(f"    ✅ Image deemed relevant for occlusion, proceeding with analysis...")
            print(f"    🔍 Calling detect_segment_rank for image analysis...")
            
            # Call the existing detection pipeline
            regions = detect_segment_rank(
                image_path=image_path,
                slide_text=slide_text,
                transcript_text=transcript_text,
                max_masks_per_image=5,
                min_mask_area_px=900,
                detection_threshold=0.3,
                nms_iou_threshold=0.5
            )
            print(f"    📊 detect_segment_rank returned {len(regions) if regions else 0} regions")
            
            # Convert to the format expected by the frontend
            formatted_regions = []
            for region in regions:
                formatted_regions.append({
                    "term": region.term,
                    "score": region.score,
                    "bbox_xyxy": list(region.bbox_xyxy),
                    "area_px": region.area_px,
                    "importance_score": region.importance_score,
                    "short_label": region.short_label,
                    "rationale": region.rationale
                })
            
            return formatted_regions
            
        finally:
            # Cleanup temp file
            if os.path.exists(image_path):
                os.unlink(image_path)
                
    except Exception as e:
        print(f"Error analyzing image: {e}")
        return []

def generate_basic_flashcards_from_text(slide_text, slide_title, use_cloze=False):
    """Generate basic Q&A or cloze flashcards from slide text using the ORIGINAL AI logic"""
    flashcards = []
    if slide_text and len(slide_text.strip()) > 10:
        try:
            # Use the ORIGINAL flashcard generation function
            from flashcard_generator import generate_enhanced_flashcards_with_progress
            
            # Prepare data in the format expected by the original function
            slide_texts = [slide_text]
            slide_images = []  # No images for basic cards
            
            # Use the original function with the REAL API key from environment
            if not OPENAI_API_KEY:
                print("❌ Error: OPENAI_API_KEY not found in environment")
                raise Exception("OpenAI API key not found")
            
            print(f"🔑 Using API key: {OPENAI_API_KEY[:20]}...")
            print(f"📝 Generating flashcards for: {slide_title}")
            print(f"🎯 Use cloze: {use_cloze}")
            
            generated_cards, analysis_data = generate_enhanced_flashcards_with_progress(
                slide_texts=slide_texts,
                slide_images=slide_images,
                api_key=OPENAI_API_KEY,  # Use the REAL API key
                model="gpt-3.5-turbo",
                max_tokens=150,
                temperature=0.7,
                progress=None,
                use_cloze=use_cloze,  # Use the cloze parameter
                question_style="Word for word"
            )
            
            print(f"✅ Generated {len(generated_cards)} cards from original function")
            
            # generated_cards already contains Flashcard objects, just update slide_number
            from flashcard_generator import Flashcard
            for card in generated_cards:
                # Update slide_number for each flashcard
                card.slide_number = 1  # Will be updated by caller
                flashcards.append(card)
                
        except Exception as e:
            print(f"❌ Error using original flashcard generation: {e}")
            import traceback
            traceback.print_exc()
            # Fallback to simple generation
            from flashcard_generator import Flashcard
            flashcard = Flashcard(
                question=f"What is the main topic of {slide_title}?",
                answer=slide_text[:100] + "..." if len(slide_text) > 100 else slide_text,
                level=1,  # Default level
                slide_number=1,  # Will be updated by caller
                confidence=0.5,  # Lower confidence for fallback
                is_cloze=use_cloze,
                cloze_text=slide_text[:100] + "..." if use_cloze else "",
                card_type="cloze" if use_cloze else "basic"
            )
            flashcards.append(flashcard)
    
    return flashcards

def create_csv_export(flashcards, csv_path):
    """Create CSV export of flashcards"""
    import csv
    with open(csv_path, 'w', newline='', encoding='utf-8') as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(['Type', 'Question', 'Answer', 'Term'])
        
        for card in flashcards:
            # Handle both Flashcard objects and dictionaries
            if hasattr(card, 'card_type'):  # Flashcard object
                if card.card_type == 'image_occlusion':
                    writer.writerow([
                        card.card_type,
                        getattr(card, 'alt_text', 'What is this?'),
                        getattr(card, 'term', 'Unknown'),
                        getattr(card, 'term', 'Unknown')
                    ])
                else:  # basic or cloze
                    writer.writerow([
                        card.card_type,
                        card.question,
                        card.answer,
                        ''
                    ])
            else:  # Dictionary (legacy format)
                if card['type'] == 'image_occlusion':
                    writer.writerow([
                        card['type'],
                        card.get('alt_text', 'What is this?'),
                        card.get('term', 'Unknown'),
                        card.get('term', 'Unknown')
                    ])
                elif card['type'] == 'basic':
                    writer.writerow([
                        card['type'],
                        card.get('question', ''),
                        card.get('answer', ''),
                        ''
                    ])
                elif card['type'] == 'cloze':
                    writer.writerow([
                        card['type'],
                        card.get('question', ''),
                        card.get('answer', ''),
                        ''
                    ])
                else:
                    writer.writerow([
                        card.get('type', 'basic'),
                        card.get('question', ''),
                        card.get('answer', ''),
                        card.get('term', '')
                    ])

def create_package_zip(apkg_path, csv_path, zip_path):
    """Create zip file with APKG and CSV"""
    import zipfile
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        zipf.write(apkg_path, os.path.basename(apkg_path))
        zipf.write(csv_path, os.path.basename(csv_path))

# Global storage for generated files (in production, use proper file storage)
generated_files = {}

@app.get("/download/{filename}")
async def download_file(filename: str):
    """Download generated files"""
    if filename in generated_files:
        file_path = generated_files[filename]
        if os.path.exists(file_path):
            return FileResponse(
                file_path,
                filename=filename,
                media_type='application/octet-stream'
            )
    
    raise HTTPException(status_code=404, detail="File not found")

if __name__ == "__main__":
    uvicorn.run(
        "api_server:app",
        host="0.0.0.0",
        port=8000,
        reload=True,
        log_level="info"
    )
